from __future__ import annotations

import base64
import logging
import re
import time
from dataclasses import dataclass
from io import BytesIO
from typing import Any

import httpx

from app.config import get_settings
from app.services.agent.models import Artifact, Source, ToolCall

logger = logging.getLogger(__name__)


def safe_filename(title: str, suffix: str) -> str:
    stem = re.sub(r"[^A-Za-z0-9._ -]+", "", title).strip().replace(" ", "-").lower()
    stem = re.sub(r"-{2,}", "-", stem)[:80].strip("-") or "agent-output"
    return f"{stem}.{suffix.lstrip('.')}"


@dataclass
class Tools:
    you_api_key: str | None = None
    tavily_api_key: str | None = None
    nimble_api_key: str | None = None
    nimble_api_endpoint: str = "https://sdk.nimbleway.com/v1/search"

    @classmethod
    def from_settings(cls) -> "Tools":
        settings = get_settings()
        return cls(
            you_api_key=settings.you_api_key,
            tavily_api_key=settings.tavily_api_key,
            nimble_api_key=settings.nimble_api_key,
            nimble_api_endpoint=settings.nimble_api_endpoint,
        )

    def search_web(self, query: str, max_results: int = 6) -> tuple[list[Source], ToolCall]:
        started = time.perf_counter()
        tool = ToolCall(name="web_search", input={"query": query, "max_results": max_results})
        try:
            if not self.you_api_key and not self.tavily_api_key and not self.nimble_api_key:
                tool.ok = False
                tool.error = "YOU_API_KEY / TAVILY_API_KEY / NIMBLE_API_KEY is not configured"
                return [], tool

            errors: list[str] = []
            if self.tavily_api_key:
                try:
                    sources = self._search_tavily(query, max_results)
                    if sources:
                        _annotate_search_sources(sources, query=query, provider="Tavily")
                        tool.output = {"provider": "Tavily", "source_count": len(sources)}
                        return sources, tool
                    errors.append("Tavily returned no results")
                except Exception as exc:
                    logger.warning("agent web_search failed: %s", exc)
                    errors.append(f"Tavily: {exc}")

            if self.you_api_key:
                try:
                    sources = self._search_you(query, max_results)
                    if sources:
                        _annotate_search_sources(sources, query=query, provider="You.com")
                        tool.output = {"provider": "You.com", "source_count": len(sources)}
                        return sources, tool
                    errors.append("You.com returned no results")
                except Exception as exc:
                    logger.warning("agent You.com web_search failed: %s", exc)
                    errors.append(f"You.com: {exc}")

            if self.nimble_api_key:
                try:
                    sources = self._search_nimble(query, max_results)
                    if sources:
                        _annotate_search_sources(sources, query=query, provider="Nimble")
                        tool.output = {"provider": "Nimble", "source_count": len(sources)}
                        return sources, tool
                    errors.append("Nimble returned no results")
                except Exception as exc:
                    logger.warning("agent Nimble web_search failed: %s", exc)
                    errors.append(f"Nimble: {exc}")

            tool.ok = False
            tool.error = "; ".join(errors) or "No web search provider returned results"
            tool.output = {"source_count": 0}
            return [], tool
        finally:
            tool.latency_ms = int((time.perf_counter() - started) * 1000)

    def _search_you(self, query: str, max_results: int) -> list[Source]:
        response = httpx.get(
            "https://ydc-index.io/v1/search",
            headers={"X-API-Key": self.you_api_key or ""},
            params={"query": query, "count": max_results},
            timeout=20,
        )
        response.raise_for_status()
        payload = response.json()
        sources: list[Source] = []
        for item in _you_result_items(payload)[:max_results]:
            if not isinstance(item, dict):
                continue
            url = str(item.get("url") or item.get("link") or "")
            snippet = _you_item_content(item)
            if url:
                sources.append(
                    Source(
                        title=str(item.get("title") or ""),
                        url=url,
                        snippet=snippet,
                    )
                )
        return sources

    def _search_tavily(self, query: str, max_results: int) -> list[Source]:
        response = httpx.post(
            "https://api.tavily.com/search",
            json={
                "api_key": self.tavily_api_key,
                "query": query,
                "max_results": max_results,
                "search_depth": "advanced",
                "include_answer": False,
                "include_raw_content": False,
            },
            timeout=20,
        )
        response.raise_for_status()
        payload = response.json()
        return [
            Source(
                title=str(item.get("title") or ""),
                url=str(item.get("url") or ""),
                snippet=str(item.get("content") or item.get("snippet") or ""),
            )
            for item in payload.get("results", [])
            if isinstance(item, dict)
        ]

    def _search_nimble(self, query: str, max_results: int) -> list[Source]:
        payload = {
            "query": query,
            "country": "US",
            "locale": "en",
            "focus": "general",
            "max_results": max_results,
            "search_depth": "lite",
            "include_answer": False,
            "output_format": "markdown",
        }
        response = httpx.post(
            self.nimble_api_endpoint,
            headers={"Content-Type": "application/json", "Authorization": _nimble_auth_header(self.nimble_api_key or "")},
            json=payload,
            timeout=20,
        )
        response.raise_for_status()
        data = response.json()
        sources: list[Source] = []
        for item in _nimble_result_items(data)[:max_results]:
            if not isinstance(item, dict):
                continue
            url = str(item.get("url") or item.get("link") or item.get("href") or "")
            snippets = item.get("snippets") or item.get("extra_snippets") or []
            if isinstance(snippets, list):
                snippet = " ".join(str(part) for part in snippets)
            else:
                snippet = str(snippets or "")
            snippet = snippet or str(item.get("description") or item.get("snippet") or item.get("content") or "")
            if url:
                sources.append(
                    Source(
                        title=str(item.get("title") or ""),
                        url=url,
                        snippet=snippet,
                    )
                )
        return sources

    def extract_urls(self, urls: list[str], max_chars_per_source: int = 2500) -> tuple[list[Source], ToolCall]:
        started = time.perf_counter()
        unique_urls = [url for index, url in enumerate(urls) if url and url not in urls[:index]][:6]
        tool = ToolCall(name="read_url", input={"urls": unique_urls})
        if not unique_urls or not self.tavily_api_key:
            tool.output = {"source_count": 0}
            tool.latency_ms = int((time.perf_counter() - started) * 1000)
            return [], tool

        try:
            response = httpx.post(
                "https://api.tavily.com/extract",
                json={"api_key": self.tavily_api_key, "urls": unique_urls},
                timeout=25,
            )
            response.raise_for_status()
            payload = response.json()
            results = payload.get("results") or []
            extracted: list[Source] = []
            for item in results:
                if not isinstance(item, dict):
                    continue
                content = str(item.get("raw_content") or item.get("content") or "")
                extracted.append(
                    Source(
                        title=str(item.get("title") or ""),
                        url=str(item.get("url") or ""),
                        content=content[:max_chars_per_source],
                    )
                )
            tool.output = {"source_count": len(extracted)}
            return extracted, tool
        except Exception as exc:
            logger.warning("agent read_url failed: %s", exc)
            tool.ok = False
            tool.error = str(exc)
            return [], tool
        finally:
            tool.latency_ms = int((time.perf_counter() - started) * 1000)

    def make_markdown_artifact(self, title: str, markdown: str) -> Artifact:
        encoded = base64.b64encode(markdown.encode("utf-8")).decode("ascii")
        return Artifact(
            kind="markdown",
            filename=safe_filename(title, "md"),
            mime_type="text/markdown",
            base64_data=encoded,
        )

    def make_docx_artifact(self, title: str, markdown: str, expected_sections: list[str] | None = None) -> tuple[Artifact, list[str]]:
        try:
            from app.services.agent.document_ast import render_docx_from_markdown

            payload, qa_issues = render_docx_from_markdown(title, markdown, expected_sections=expected_sections)
            mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            kind = "docx"
            filename = safe_filename(title, "docx")
            issue_codes = [issue.code for issue in qa_issues]
        except Exception as exc:
            logger.warning("agent docx artifact failed; returning markdown: %s", exc)
            payload = markdown.encode("utf-8")
            mime = "text/markdown"
            kind = "markdown"
            filename = safe_filename(title, "md")
            issue_codes = ["docx_render_failed"]
        artifact = Artifact(
            kind=kind,  # type: ignore[arg-type]
            filename=filename,
            mime_type=mime,
            base64_data=base64.b64encode(payload).decode("ascii"),
        )
        return artifact, issue_codes

    def make_pptx_artifact(
        self,
        title: str,
        markdown: str,
        expected_slides: list[str] | None = None,
        template_id: str | None = None,
        user_id: str | None = None,
        render_plan: dict[str, Any] | None = None,
        design_system_id: str | None = None,
        repair_actions: list[dict[str, Any]] | None = None,
    ) -> tuple[Artifact, dict[str, Any]]:
        try:
            from app.services.agent.pptx_design import render_agentdeck_pptx_from_markdown, render_agentdeck_pptx_from_render_plan
            from app.services.pptx_render_qa import run_pptx_render_qa
            from app.services.components.render_plan import PptxRenderPlan

            if render_plan:
                plan = PptxRenderPlan.model_validate(render_plan)
                design = render_agentdeck_pptx_from_render_plan(
                    plan,
                    design_system_id=design_system_id or "agentdeck_v1",
                    repair_actions=repair_actions,
                )
                source = "structured_deck_plan"
            else:
                design = render_agentdeck_pptx_from_markdown(title=title, markdown=markdown, template_id=template_id, user_id=user_id)
                source = "markdown_bridge"
            payload = design.payload
            qa = run_pptx_render_qa(payload)
            issue_codes = [str(issue.get("type") or "unknown") for issue in qa.get("issues", []) if isinstance(issue, dict)]
            mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            kind = "pptx"
            filename = safe_filename(title, "pptx")
            metadata: dict[str, Any] = {
                "design_system": design.design_system_id,
                "theme": design.theme,
                "template_id": template_id,
                "template_applied": bool(template_id and design.design_system_id != "agentdeck_v1"),
                "deck_source": source,
                "qa_available": bool(qa.get("available")),
                "qa_issue_codes": issue_codes,
                "slide_count": qa.get("slide_count") or design.slide_count,
                "layout_counts": design.layout_counts,
                "design_ledger": design.design_ledger,
                "repair_actions": design.repair_actions,
                "expected_slide_count": len(expected_slides or []),
            }
        except Exception as exc:
            logger.warning("agent pptx artifact failed; returning markdown: %s", exc)
            payload = markdown.encode("utf-8")
            mime = "text/markdown"
            kind = "markdown"
            filename = safe_filename(title, "md")
            metadata = {
                "qa_available": False,
                "qa_issue_codes": ["pptx_render_failed"],
                "error": str(exc),
                "template_id": template_id,
                "template_applied": False,
                "deck_source": "structured_deck_plan" if render_plan else "markdown_bridge",
                "expected_slide_count": len(expected_slides or []),
            }
        artifact = Artifact(
            kind=kind,  # type: ignore[arg-type]
            filename=filename,
            mime_type=mime,
            base64_data=base64.b64encode(payload).decode("ascii"),
        )
        return artifact, metadata


def render_pptx_from_markdown(title: str, markdown: str) -> bytes:
    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    deck = Presentation()
    deck.slide_width = Inches(13.333)
    deck.slide_height = Inches(7.5)
    slides = _slides_from_markdown(title, markdown)
    for index, slide_spec in enumerate(slides):
        layout = deck.slide_layouts[0] if index == 0 else deck.slide_layouts[1]
        slide = deck.slides.add_slide(layout)
        slide.shapes.title.text = slide_spec["title"]
        if index == 0:
            subtitle = slide.placeholders[1] if len(slide.placeholders) > 1 else None
            if subtitle is not None:
                subtitle.text = slide_spec.get("subtitle") or "Prepared by Fronei"
        else:
            body = slide.placeholders[1]
            body.text = ""
            frame = body.text_frame
            frame.clear()
            bullets = slide_spec.get("bullets") or ["No visible slide bullets were generated."]
            for bullet_index, bullet in enumerate(bullets[:6]):
                paragraph = frame.paragraphs[0] if bullet_index == 0 else frame.add_paragraph()
                paragraph.text = bullet
                paragraph.level = 0
                paragraph.font.size = Pt(22)
        notes = slide.notes_slide.notes_text_frame
        notes.text = slide_spec.get("notes") or ""
        if index == 0:
            slide.shapes.title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    package = BytesIO()
    deck.save(package)
    return package.getvalue()


def _slides_from_markdown(title: str, markdown: str) -> list[dict[str, Any]]:
    deck_title = title.strip() or "Fronei presentation"
    slides: list[dict[str, Any]] = [{"title": deck_title, "subtitle": "Fronei slide deck", "bullets": [], "notes": ""}]
    current: dict[str, Any] | None = None
    in_notes = False
    for raw_line in (markdown or "").splitlines():
        line = raw_line.strip()
        if not line:
            continue
        if line.startswith("# "):
            deck_title = line[2:].strip() or deck_title
            slides[0]["title"] = deck_title
            continue
        if line.startswith("## "):
            current = {"title": line[3:].strip() or "Slide", "bullets": [], "notes": ""}
            slides.append(current)
            in_notes = False
            continue
        if current is None:
            continue
        lowered = line.lower()
        if lowered.startswith(("notes:", "speaker notes:", "presenter notes:")):
            current["notes"] = line.split(":", 1)[1].strip()
            in_notes = True
            continue
        if line.startswith(("- ", "* ")):
            current["bullets"].append(_clean_slide_text(line[2:]))
            in_notes = False
            continue
        if re.match(r"^\d+[.)]\s+", line):
            current["bullets"].append(_clean_slide_text(re.sub(r"^\d+[.)]\s+", "", line)))
            in_notes = False
            continue
        if in_notes:
            current["notes"] = (current.get("notes", "") + " " + line).strip()
        elif len(current["bullets"]) < 4:
            current["bullets"].append(_clean_slide_text(line))
        else:
            current["notes"] = (current.get("notes", "") + " " + line).strip()
    if len(slides) == 1:
        slides.append({"title": "Overview", "bullets": [_clean_slide_text(markdown or "Presentation content")], "notes": ""})
    return slides[:24]


def _clean_slide_text(value: str) -> str:
    text = re.sub(r"\*\*([^*]+)\*\*", r"\1", value or "")
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
    return text.strip()


_DEFAULT_SOURCE_CONTEXT_MAX_CHARS = 2500


def source_context(sources: list[Source], *, max_chars_per_source: int | None = None) -> str:
    cap = max_chars_per_source or _DEFAULT_SOURCE_CONTEXT_MAX_CHARS
    lines: list[str] = []
    for idx, source in enumerate(sources, start=1):
        body = source.content or source.snippet
        lines.append(f"[S{idx}] {source.title}\nURL: {source.url}\n{body[:cap]}")
    return "\n\n".join(lines)


def _nimble_auth_header(api_key: str) -> str:
    if api_key.lower().startswith(("bearer ", "basic ")):
        return api_key
    return f"Bearer {api_key}"


def _annotate_search_sources(sources: list[Source], *, query: str, provider: str) -> None:
    for source in sources:
        source.query = query
        source.provider = provider


def _you_result_items(data: dict) -> list[dict]:
    results = data.get("results")
    if isinstance(results, dict):
        items: list[dict] = []
        for key in ("web", "news"):
            value = results.get(key)
            if isinstance(value, list):
                items.extend(item for item in value if isinstance(item, dict))
        return items
    if isinstance(results, list):
        return [item for item in results if isinstance(item, dict)]
    return []


def _you_item_content(item: dict) -> str:
    snippets = item.get("snippets") or item.get("highlights") or []
    if isinstance(snippets, list):
        snippet = " ".join(str(part) for part in snippets)
    else:
        snippet = str(snippets or "")
    if snippet:
        return snippet
    contents = item.get("contents") if isinstance(item.get("contents"), dict) else {}
    return str(
        item.get("description")
        or item.get("snippet")
        or item.get("content")
        or contents.get("markdown")
        or contents.get("html")
        or ""
    )


def _nimble_result_items(data: dict) -> list[dict]:
    parsing = data.get("parsing") if isinstance(data.get("parsing"), dict) else {}
    entities = parsing.get("entities") if isinstance(parsing.get("entities"), dict) else {}
    for key in ("SearchResult", "OrganicResult", "organic_results", "search_results", "results"):
        value = entities.get(key)
        if isinstance(value, list):
            return value
    for key in ("organic_results", "search_results", "results", "items"):
        value = data.get(key)
        if isinstance(value, list):
            return value
    return []
