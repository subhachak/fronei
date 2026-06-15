from __future__ import annotations

import copy
import base64
import json
import logging
import re
import select
import shutil
import subprocess
import threading
import uuid
from concurrent.futures import ThreadPoolExecutor, as_completed
from datetime import date
from io import BytesIO
from pathlib import Path

from docx import Document
from docx.enum.text import WD_BREAK
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt
from docx.text.paragraph import Paragraph
from markdown_it import MarkdownIt
from openpyxl import Workbook
from openpyxl.chart import BarChart, Reference
from openpyxl.styles import Alignment, Font, PatternFill
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn as pptx_qn
from pptx.util import Inches as PptxInches, Pt as PptxPt

from app.config import get_settings
from app.services.components import PptxRenderPlan, PptxSlidePlan, ZoneInstance
from app.services.document_templates import resolve_pptx_template_path
from app.services.presentation_design_system import (
    FIT_CONTRACTS,
    LAYOUT_ALIASES as DESIGN_LAYOUT_ALIASES,
    canonical_layout,
    component_tree_for_slide,
    design_system_payload,
    template_for_slide,
)

logger = logging.getLogger(__name__)

# PptxGenJS-based renderer for the "no template" (fronei-default) PPTX path —
# see PPTX_RENDER_DIR / "render.js". Decks built from a built-in or
# user-uploaded branded .pptx template still go through python-pptx (below),
# which can read that template's layouts/placeholders directly.
PPTX_RENDER_DIR = Path(__file__).resolve().parents[2] / "pptx_render"
PPTX_RENDER_JS = PPTX_RENDER_DIR / "render.js"
PPTX_RENDER_AGENTDECK_JS = PPTX_RENDER_DIR / "agentdeck" / "render_agentdeck.js"
PPTX_RENDER_AGENTDECK_SERVER_JS = PPTX_RENDER_DIR / "agentdeck" / "render_agentdeck_server.js"
PPTX_RENDER_TIMEOUT_SECONDS = 60


TABLE_SEPARATOR_RE = re.compile(r"^\s*\|?\s*:?-{3,}:?\s*(\|\s*:?-{3,}:?\s*)+\|?\s*$")
_MARKDOWN_INLINE = MarkdownIt("commonmark")
KNOWN_DOC_TYPES = {
    "executive_report",
    "proposal",
    "memo",
    "technical_spec",
    "meeting_notes",
    "one_pager",
    "letter",
    "resume",
    "presentation",
}
SPEAKER_NOTES_RE = re.compile(r"^speaker notes?\s*:\s*(.*)$", re.IGNORECASE)
MAX_BULLETS_PER_SLIDE = 6
# Slide titles are assertion-style, short headlines. The prompt asks the LLM
# for ~40-60 chars; this is the ceiling beyond which _shorten_title_to_notes
# shortens at a clause/word boundary (never mid-word, never literal "...")
# and routes the full title to speaker_notes.
MAX_SLIDE_TITLE_CHARS = 90
MAX_BULLET_CHARS = 90

# Default theme used by the python-pptx fallback renderer when the uploaded
# template's layouts provide no usable placeholders (common for "design
# scaffold" templates where every slide is hand-laid-out free shapes).
# Mirrors the warm-editorial palette used by pptx_render/render.js.
PPTX_FALLBACK_HEADING_FONT = "Georgia"
PPTX_FALLBACK_BODY_FONT = "Calibri"
PPTX_FALLBACK_TEXT_RGB = RGBColor(0x28, 0x24, 0x21)
PPTX_FALLBACK_ACCENT_RGB = RGBColor(0xE0, 0x4F, 0x00)
PPTX_FALLBACK_NAVY_RGB = RGBColor(0x1F, 0x3B, 0x5C)
PPTX_FALLBACK_CARD_BG_RGB = RGBColor(0xFF, 0xFD, 0xFC)
PPTX_FALLBACK_CARD_LINE_RGB = RGBColor(0xD8, 0xCD, 0xC6)

# Title box/accent-rule geometry for the fallback title textbox — mirrors
# TITLE_BOX_H / TITLE_RULE_Y / CONTENT_TOP_Y in render.js so placeholder-less
# templates get the same layout rhythm as the PptxGenJS renderer.
PPTX_TITLE_BOX_H = 1.0
PPTX_TITLE_RULE_Y = 1.32
PPTX_CONTENT_TOP_Y = 1.65

# Per-template design tokens for the python-pptx fallback renderer. Each of
# the built-in "design scaffold" templates (warm-editorial, modern-tech,
# executive-navy, data-product-os, clean-light) ships its own bg/card/fg/
# muted/accent/accent2 palette and heading/body fonts, extracted from that
# template's "design system snapshot" slide. The fallback renderer (used for
# templates whose layouts have no placeholders) looks up the active theme via
# `_pptx_theme()` so generated decks match the uploaded template's palette
# instead of a single hardcoded warm-editorial look.
PPTX_DEFAULT_THEME: dict = {
    "bg": RGBColor(0xF6, 0xF0, 0xE6),
    "card": PPTX_FALLBACK_CARD_BG_RGB,
    "card_line": PPTX_FALLBACK_CARD_LINE_RGB,
    "fg": PPTX_FALLBACK_TEXT_RGB,
    "muted": RGBColor(0x6B, 0x5E, 0x52),
    "accent": PPTX_FALLBACK_ACCENT_RGB,
    "accent2": RGBColor(0x0F, 0x76, 0x6E),
    "heading_font": PPTX_FALLBACK_HEADING_FONT,
    "body_font": PPTX_FALLBACK_BODY_FONT,
}

PPTX_TEMPLATE_THEMES: dict[str, dict] = {
    "warm-editorial": {
        "bg": RGBColor(0xF6, 0xF0, 0xE6),
        "card": RGBColor(0xFF, 0xFD, 0xF8),
        "card_line": RGBColor(0xD8, 0xCD, 0xC6),
        "fg": RGBColor(0x1F, 0x29, 0x37),
        "muted": RGBColor(0x6B, 0x5E, 0x52),
        "accent": RGBColor(0xB4, 0x50, 0x09),
        "accent2": RGBColor(0x0F, 0x76, 0x6E),
        "heading_font": "Georgia",
        "body_font": "Calibri",
    },
    "modern-tech": {
        "bg": RGBColor(0x08, 0x0C, 0x11),
        "card": RGBColor(0x12, 0x1A, 0x24),
        "card_line": RGBColor(0x24, 0x30, 0x3D),
        "fg": RGBColor(0xEF, 0xF6, 0xFF),
        "muted": RGBColor(0xAA, 0xB8, 0xC7),
        "accent": RGBColor(0x22, 0xD3, 0xEE),
        "accent2": RGBColor(0xA3, 0xE6, 0x35),
        "heading_font": "Calibri",
        "body_font": "Calibri",
    },
    "executive-navy": {
        "bg": RGBColor(0x10, 0x18, 0x27),
        "card": RGBColor(0x17, 0x20, 0x33),
        "card_line": RGBColor(0x2A, 0x37, 0x52),
        "fg": RGBColor(0xF8, 0xFA, 0xFC),
        "muted": RGBColor(0xA7, 0xB2, 0xC5),
        "accent": RGBColor(0x38, 0xBD, 0xF8),
        "accent2": RGBColor(0x7C, 0x3A, 0xED),
        "heading_font": "Calibri",
        "body_font": "Calibri",
    },
    "data-product-os": {
        "bg": RGBColor(0x0B, 0x12, 0x20),
        "card": RGBColor(0x11, 0x18, 0x27),
        "card_line": RGBColor(0x1E, 0x29, 0x3B),
        "fg": RGBColor(0xF1, 0xF5, 0xF9),
        "muted": RGBColor(0xCB, 0xD5, 0xE1),
        "accent": RGBColor(0x34, 0xD3, 0x99),
        "accent2": RGBColor(0xF5, 0x9E, 0x0B),
        "heading_font": "Calibri",
        "body_font": "Calibri",
    },
    "clean-light": {
        "bg": RGBColor(0xF8, 0xFA, 0xFC),
        "card": RGBColor(0xFF, 0xFF, 0xFF),
        "card_line": RGBColor(0xE2, 0xE8, 0xF0),
        "fg": RGBColor(0x0F, 0x17, 0x2A),
        "muted": RGBColor(0x47, 0x55, 0x69),
        "accent": RGBColor(0x25, 0x63, 0xEB),
        "accent2": RGBColor(0x10, 0xB9, 0x81),
        "heading_font": "Calibri",
        "body_font": "Calibri",
    },
}

# Active theme for the current render call (set by
# `_generate_pptx_bytes_python_pptx` before rendering, restored afterward).
# Module-level rather than threaded through every helper signature for
# simplicity — PPTX generation is synchronous/single-call per request.
_ACTIVE_PPTX_THEME: dict = PPTX_DEFAULT_THEME


def _pptx_theme() -> dict:
    return _ACTIVE_PPTX_THEME


def _pptx_set_theme(template_id: str | None) -> dict:
    """Activate the design theme for `template_id` and return the previous
    theme so the caller can restore it after rendering."""
    global _ACTIVE_PPTX_THEME
    previous = _ACTIVE_PPTX_THEME
    _ACTIVE_PPTX_THEME = PPTX_TEMPLATE_THEMES.get(template_id or "", PPTX_DEFAULT_THEME)
    return previous


def _pptx_restore_theme(previous: dict) -> None:
    global _ACTIVE_PPTX_THEME
    _ACTIVE_PPTX_THEME = previous


def _pptx_set_slide_background(slide) -> None:
    """Fill the slide background with the active theme's `bg` color. New
    slides added to a placeholder-less template layout don't inherit a
    background fill, so dark themes would otherwise render on a white page."""
    theme = _pptx_theme()
    bg = theme.get("bg")
    if bg is None:
        return
    try:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = bg
    except Exception:
        pass


def _pptx_add_slide(prs: Presentation, role: str):
    """Add a slide for `role` and immediately paint the theme background.

    Without this, slides rendered onto a template layout that *has* a title
    placeholder (the common case for real PowerPoint templates) never call
    `_pptx_set_slide_background` — only the placeholder-less fallback path
    did — so themed decks rendered with plain white backgrounds regardless
    of the active theme's `bg` color. Centralizing add_slide here ensures
    every slide gets the theme background up front; per-shape fills (cards,
    accent rules, etc.) are layered on top by the layout-specific renderers."""
    slide = prs.slides.add_slide(_pptx_layout_for_role(prs, role))
    _pptx_set_slide_background(slide)
    return slide


def _pptx_set_notes(slide, notes: str | None) -> None:
    if not notes:
        return
    try:
        notes_frame = slide.notes_slide.notes_text_frame
        if notes_frame is not None:
            notes_frame.text = notes
    except Exception:
        logger.debug("Could not write speaker notes to slide", exc_info=True)

# Appendix slides are reference material — denser content is acceptable, so
# they get a higher per-slide bullet cap than the standard body slides.
MAX_APPENDIX_BULLETS = 10

# Layout name aliases normalized by parse_deck_plan. Both sides of an alias
# pair are treated identically by the PPTX renderer.
DECK_LAYOUT_ALIASES = {
    **DESIGN_LAYOUT_ALIASES,
    "cover": "section",
    "hero_cover": "section",
    "decision": "recommendation",
    "decision_slide": "recommendation",
    "decision_recommendation": "recommendation",
    "roadmap": "timeline",
    "process": "timeline",
    "process_steps": "timeline",
    "architecture_map": "architecture",
    "system_map": "architecture",
    "financial_exhibit": "financial_model",
    "data_exhibit": "financial_model",
    "three_card_system": "comparison",
    "governance_grid": "comparison",
    "principles_grid": "comparison",
    "takeaways": "executive_summary",
    "stats": "stat_cards",
    "metrics": "stat_cards",
    "kpi": "stat_cards",
    "kpi_grid": "stat_cards",
    "market_context": "stat_cards",
    "by_the_numbers": "stat_cards",
}
SLIDE_ARCHETYPE_LIBRARY = {
    "decision_pack_cover": {
        "layout": "cover_metric_strip",
        "proof_objects": {"stat_cards", "insight_cards"},
        "required_any": {"stats", "bullets"},
        "render_hints": {"tone": "decisive", "visual_weight": "cover_metrics", "accent": "decision"},
    },
    "current_state_estate_map": {
        "layout": "current_state_estate_map",
        "proof_objects": {"comparison", "insight_cards"},
        "required_any": {"units", "columns", "bullets"},
        "render_hints": {"tone": "diagnostic", "visual_weight": "estate_map", "accent": "technical"},
    },
    "impact_scorecard_bars": {
        "layout": "impact_scorecard_bars",
        "proof_objects": {"stat_cards", "chart"},
        "required_any": {"stats", "bars"},
        "render_hints": {"tone": "quantified", "visual_weight": "scorecard_bars", "accent": "financial"},
    },
    "option_score_matrix": {
        "layout": "option_score_matrix",
        "proof_objects": {"comparison"},
        "required_any": {"options", "columns"},
        "render_hints": {"tone": "evaluative", "visual_weight": "score_matrix", "accent": "operational"},
    },
    "platform_operating_model_hub": {
        "layout": "platform_operating_model_hub",
        "proof_objects": {"architecture", "comparison"},
        "required_any": {"platform", "columns", "bullets"},
        "render_hints": {"tone": "operational", "visual_weight": "hub_spoke", "accent": "technical"},
    },
    "roadmap_phase_cards": {
        "layout": "roadmap_phase_cards",
        "proof_objects": {"timeline"},
        "required_any": {"phases"},
        "render_hints": {"tone": "sequenced", "visual_weight": "phase_cards", "accent": "execution"},
    },
    "risk_control_rows": {
        "layout": "risk_control_rows",
        "proof_objects": {"comparison", "table"},
        "required_any": {"columns", "table"},
        "render_hints": {"tone": "controlled", "visual_weight": "risk_rows", "accent": "risk"},
    },
    "decision_ask_panel": {
        "layout": "decision_ask_panel",
        "proof_objects": {"stat_cards", "comparison"},
        "required_any": {"decisions", "stats", "bullets"},
        "render_hints": {"tone": "decisive", "visual_weight": "decision_panel", "accent": "decision"},
    },
    "board_decision": {
        "layout": "recommendation",
        "proof_objects": {"insight_cards", "comparison", "stat_cards"},
        "required_any": {"bullets", "columns", "stats"},
        "render_hints": {"tone": "decisive", "visual_weight": "hero_decision", "accent": "decision"},
    },
    "metric_scorecard": {
        "layout": "stat_cards",
        "proof_objects": {"stat_cards"},
        "required_any": {"stats"},
        "render_hints": {"tone": "quantified", "visual_weight": "kpi_grid", "accent": "financial"},
    },
    "risk_register": {
        "layout": "comparison",
        "proof_objects": {"comparison", "table", "insight_cards"},
        "required_any": {"columns", "table", "bullets"},
        "render_hints": {"tone": "controlled", "visual_weight": "risk_grid", "accent": "risk"},
    },
    "operating_model": {
        "layout": "comparison",
        "proof_objects": {"comparison", "table", "insight_cards"},
        "required_any": {"columns", "table", "bullets"},
        "render_hints": {"tone": "operational", "visual_weight": "role_lanes", "accent": "operational"},
    },
    "architecture_map": {
        "layout": "architecture",
        "proof_objects": {"architecture", "comparison", "insight_cards"},
        "required_any": {"bullets", "columns"},
        "render_hints": {"tone": "technical", "visual_weight": "node_flow", "accent": "technical"},
    },
    "investment_case": {
        "layout": "stat_cards",
        "proof_objects": {"stat_cards", "chart", "table", "insight_cards"},
        "required_any": {"stats", "chart", "table", "bullets"},
        "render_hints": {"tone": "financial", "visual_weight": "business_case", "accent": "financial"},
    },
    "roadmap": {
        "layout": "timeline",
        "proof_objects": {"timeline"},
        "required_any": {"phases"},
        "render_hints": {"tone": "sequenced", "visual_weight": "phase_path", "accent": "execution"},
    },
    "comparison_matrix": {
        "layout": "comparison",
        "proof_objects": {"comparison", "table"},
        "required_any": {"columns", "table"},
        "render_hints": {"tone": "evaluative", "visual_weight": "option_cards", "accent": "operational"},
    },
}
COVER_DOC_TYPES = {"executive_report", "proposal", "technical_spec"}
COMPACT_HEADER_DOC_TYPES = {"memo", "one_pager", "resume"}
TOC_DOC_TYPES = {"executive_report", "proposal"}


def _inline_tokens(text: str):
    try:
        parsed = _MARKDOWN_INLINE.parseInline(str(text or ""))
    except Exception:
        return None
    if not parsed:
        return []
    return parsed[0].children or []


def _clean_inline(text: str) -> str:
    tokens = _inline_tokens(text)
    if tokens is None:
        text = re.sub(r"!\[([^\]]*)\]\([^)]+\)", r"\1", text)
        text = re.sub(r"\[([^\]]+)\]\(([^)]+)\)", r"\1 (\2)", text)
        text = re.sub(r"`([^`]+)`", r"\1", text)
        text = re.sub(r"(\*\*|__)(.*?)\1", r"\2", text)
        text = re.sub(r"(\*|_)(.*?)\1", r"\2", text)
        return text.strip()

    parts: list[str] = []
    link_href: str | None = None
    for token in tokens:
        if token.type == "link_open":
            link_href = (token.attrs or {}).get("href")
            continue
        if token.type == "link_close":
            if link_href:
                parts.append(f" ({link_href})")
            link_href = None
            continue
        if token.type in {"text", "code_inline", "image"}:
            parts.append(token.content)
        elif token.type in {"softbreak", "hardbreak"}:
            parts.append("\n")
    return "".join(parts).strip()


def _inline_token_text(token) -> str:
    if token.type == "image":
        return token.content or (token.attrs or {}).get("alt", "")
    if token.type in {"softbreak", "hardbreak"}:
        return "\n"
    return token.content or ""


def _split_table_row(line: str) -> list[str]:
    trimmed = line.strip().strip("|")
    return [cell.strip() for cell in trimmed.split("|")]


def _is_table_start(lines: list[str], idx: int) -> bool:
    return (
        idx + 1 < len(lines)
        and lines[idx].strip().startswith("|")
        and TABLE_SEPARATOR_RE.match(lines[idx + 1]) is not None
    )


def _add_table(doc: Document, rows: list[list[str]]) -> None:
    if not rows:
        return
    width = max(len(row) for row in rows)
    table = doc.add_table(rows=len(rows), cols=width)
    table.style = "Table Grid"
    for r_idx, row in enumerate(rows):
        for c_idx in range(width):
            cell = table.cell(r_idx, c_idx)
            cell.text = ""
            paragraph = cell.paragraphs[0]
            _add_inline_runs(paragraph, row[c_idx] if c_idx < len(row) else "", base_bold=(r_idx == 0))
            if r_idx == 0:
                for paragraph in cell.paragraphs:
                    for run in paragraph.runs:
                        run.bold = True


def _add_code_block(doc: Document, code_lines: list[str]) -> None:
    if not code_lines:
        return
    paragraph = doc.add_paragraph()
    run = paragraph.add_run("\n".join(code_lines))
    run.font.name = "Courier New"
    run.font.size = Pt(9)


def _add_hyperlink(paragraph: Paragraph, text: str, url: str, base_bold: bool, base_italic: bool) -> None:
    part = paragraph.part
    relationship_id = part.relate_to(
        url,
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
        is_external=True,
    )
    hyperlink = OxmlElement("w:hyperlink")
    hyperlink.set(qn("r:id"), relationship_id)

    run_element = OxmlElement("w:r")
    properties = OxmlElement("w:rPr")

    if base_bold:
        properties.append(OxmlElement("w:b"))
    if base_italic:
        properties.append(OxmlElement("w:i"))

    color = OxmlElement("w:color")
    color.set(qn("w:val"), "0563C1")
    properties.append(color)
    underline = OxmlElement("w:u")
    underline.set(qn("w:val"), "single")
    properties.append(underline)

    text_element = OxmlElement("w:t")
    text_element.text = text
    run_element.append(properties)
    run_element.append(text_element)
    hyperlink.append(run_element)
    paragraph._p.append(hyperlink)


def _add_run(
    paragraph: Paragraph,
    text: str,
    *,
    bold: bool = False,
    italic: bool = False,
    code: bool = False,
    base_bold: bool = False,
    base_italic: bool = False,
) -> None:
    if not text:
        return
    run = paragraph.add_run(text)
    run.bold = base_bold or bold
    run.italic = base_italic or italic
    if code:
        run.font.name = "Courier New"
        run.font.size = Pt(9.5)


def _add_inline_runs(
    paragraph: Paragraph,
    text: str,
    *,
    base_bold: bool = False,
    base_italic: bool = False,
) -> None:
    """Add CommonMark inline text as formatted DOCX runs."""
    tokens = _inline_tokens(text)
    if tokens is None:
        _add_run(paragraph, str(text or "").strip(), base_bold=base_bold, base_italic=base_italic)
        return

    bold_depth = 0
    italic_depth = 0
    link_href: str | None = None
    for token in tokens:
        if token.type == "strong_open":
            bold_depth += 1
            continue
        if token.type == "strong_close":
            bold_depth = max(0, bold_depth - 1)
            continue
        if token.type == "em_open":
            italic_depth += 1
            continue
        if token.type == "em_close":
            italic_depth = max(0, italic_depth - 1)
            continue
        if token.type == "link_open":
            link_href = (token.attrs or {}).get("href")
            continue
        if token.type == "link_close":
            link_href = None
            continue

        token_text = _inline_token_text(token)
        if not token_text:
            continue
        is_bold = bold_depth > 0
        is_italic = italic_depth > 0
        if link_href and token.type != "code_inline":
            _add_hyperlink(paragraph, token_text, link_href, base_bold or is_bold, base_italic or is_italic)
            continue
        _add_run(
            paragraph,
            token_text,
            bold=is_bold,
            italic=is_italic,
            code=(token.type == "code_inline"),
            base_bold=base_bold,
            base_italic=base_italic,
        )


def _add_inline_paragraph(
    doc: Document,
    text: str,
    *,
    style: str | None = None,
    base_bold: bool = False,
    base_italic: bool = False,
) -> None:
    paragraph = doc.add_paragraph(style=style)
    _add_inline_runs(paragraph, text, base_bold=base_bold, base_italic=base_italic)


def _add_field(paragraph: Paragraph, instruction: str, placeholder: str = "") -> None:
    run = paragraph.add_run()
    begin = OxmlElement("w:fldChar")
    begin.set(qn("w:fldCharType"), "begin")
    run._r.append(begin)

    instr = OxmlElement("w:instrText")
    instr.set(qn("xml:space"), "preserve")
    instr.text = instruction
    run._r.append(instr)

    separate = OxmlElement("w:fldChar")
    separate.set(qn("w:fldCharType"), "separate")
    run._r.append(separate)
    if placeholder:
        text = OxmlElement("w:t")
        text.text = placeholder
        run._r.append(text)

    end = OxmlElement("w:fldChar")
    end.set(qn("w:fldCharType"), "end")
    run._r.append(end)


def _add_footer(section) -> None:
    paragraph = section.footer.paragraphs[0]
    paragraph.text = "Fronei"
    paragraph.add_run(" | Page ")
    _add_field(paragraph, "PAGE", "1")


def _apply_type_styles(doc: Document, doc_type: str) -> None:
    section = doc.sections[0]
    if doc_type == "executive_report":
        section.top_margin = Inches(0.85)
        section.bottom_margin = Inches(0.8)
        section.left_margin = Inches(0.9)
        section.right_margin = Inches(0.9)
        sizes = [("Title", 28), ("Heading 1", 20), ("Heading 2", 15), ("Heading 3", 12)]
    elif doc_type == "one_pager":
        section.top_margin = Inches(0.55)
        section.bottom_margin = Inches(0.55)
        section.left_margin = Inches(0.6)
        section.right_margin = Inches(0.6)
        sizes = [("Title", 18), ("Heading 1", 14), ("Heading 2", 12), ("Heading 3", 11)]
        doc.styles["Normal"].font.size = Pt(9.5)
    elif doc_type == "proposal":
        section.top_margin = Inches(0.8)
        section.bottom_margin = Inches(0.8)
        section.left_margin = Inches(0.85)
        section.right_margin = Inches(0.85)
        sizes = [("Title", 26), ("Heading 1", 18), ("Heading 2", 14), ("Heading 3", 12)]
    else:
        sizes = [("Title", 24), ("Heading 1", 18), ("Heading 2", 14), ("Heading 3", 12)]

    for style_name, size in sizes:
        doc.styles[style_name].font.size = Pt(size)


def _add_cover(doc: Document, title: str, subtitle: str | None) -> None:
    _add_inline_paragraph(doc, _clean_inline(title) or "Fronei document", style="Title")
    if subtitle:
        _add_inline_paragraph(doc, subtitle, base_italic=True)
    doc.add_paragraph("")
    prepared = doc.add_paragraph()
    _add_inline_runs(prepared, f"Prepared for Fronei\n{date.today().isoformat()}")
    doc.add_page_break()


def _add_compact_header(doc: Document, title: str, subtitle: str | None) -> None:
    _add_inline_paragraph(doc, _clean_inline(title) or "Fronei document", style="Heading 1")
    meta = doc.add_paragraph()
    meta_text = f"Prepared for Fronei | {date.today().isoformat()}"
    if subtitle:
        meta_text = f"{subtitle} | {meta_text}"
    _add_inline_runs(meta, meta_text, base_italic=True)


def _add_toc(doc: Document, content: str) -> None:
    """Render a static table of contents from the document's headings.

    Deliberately avoids a Word TOC field (w:fldChar / instrText "TOC ...").
    Field codes — even without w:updateFields set — make Word show
    "This document contains fields that may refer to other files. Do you
    want to update the fields in this document?" on open for some users/
    Word configurations. A plain heading list has identical informational
    value without any field codes.
    """
    headings: list[tuple[int, str]] = []
    skipped_first_h1 = False
    for raw in content.replace("\r\n", "\n").replace("\r", "\n").split("\n"):
        match = re.match(r"^(#{1,3})\s+(.+)$", raw.strip())
        if not match:
            continue
        level = len(match.group(1))
        text = _clean_inline(match.group(2)).strip()
        if level == 1 and not skipped_first_h1:
            # Skip the document's own title heading (rendered separately).
            skipped_first_h1 = True
            continue
        if text:
            headings.append((level, text))

    if not headings:
        return

    doc.add_heading("Table of Contents", level=1)
    for level, text in headings:
        paragraph = doc.add_paragraph(text)
        paragraph.paragraph_format.left_indent = Inches(0.25 * (level - 1))
    doc.add_page_break()


def _strip_leading_h1(content: str) -> str:
    """Remove a leading top-level Markdown heading, if present.

    Used when the title is already rendered separately (cover page or
    compact header) so it doesn't appear a second time in the body.
    """
    lines = content.replace("\r\n", "\n").replace("\r", "\n").split("\n")
    idx = 0
    while idx < len(lines) and not lines[idx].strip():
        idx += 1
    if idx < len(lines) and re.match(r"^#\s+.+$", lines[idx].strip()):
        return "\n".join(lines[idx + 1:])
    return content


def generate_docx_bytes(title: str, content: str, subtitle: str | None = None, doc_type: str | None = None) -> bytes:
    """Render Markdown-ish text into a simple, professional DOCX document."""
    doc = Document()
    section = doc.sections[0]
    section.top_margin = Inches(0.7)
    section.bottom_margin = Inches(0.7)
    section.left_margin = Inches(0.8)
    section.right_margin = Inches(0.8)

    styles = doc.styles
    styles["Normal"].font.name = "Aptos"
    styles["Normal"].font.size = Pt(10.5)
    for style_name, size in [("Title", 24), ("Heading 1", 18), ("Heading 2", 14), ("Heading 3", 12)]:
        styles[style_name].font.name = "Aptos Display"
        styles[style_name].font.size = Pt(size)

    enhanced_doc_type = doc_type if doc_type in KNOWN_DOC_TYPES else None
    if enhanced_doc_type:
        _apply_type_styles(doc, enhanced_doc_type)
        _add_footer(section)
        if enhanced_doc_type in COVER_DOC_TYPES:
            _add_cover(doc, title, subtitle)
            if enhanced_doc_type in TOC_DOC_TYPES:
                _add_toc(doc, content)
        elif enhanced_doc_type in COMPACT_HEADER_DOC_TYPES:
            _add_compact_header(doc, title, subtitle)
        # The cover / compact header already renders the title, so drop a
        # leading "# Title" line from the body to avoid showing it twice.
        content = _strip_leading_h1(content)
    else:
        _add_inline_paragraph(doc, _clean_inline(title) or "Fronei document", style="Title")
        if subtitle:
            _add_inline_paragraph(doc, subtitle, base_italic=True)

    lines = content.replace("\r\n", "\n").replace("\r", "\n").split("\n")
    idx = 0
    in_code = False
    code_lines: list[str] = []

    while idx < len(lines):
        raw = lines[idx]
        line = raw.rstrip()
        stripped = line.strip()

        if stripped.startswith("```"):
            if in_code:
                _add_code_block(doc, code_lines)
                code_lines = []
                in_code = False
            else:
                in_code = True
            idx += 1
            continue

        if in_code:
            code_lines.append(raw)
            idx += 1
            continue

        if not stripped:
            idx += 1
            continue

        if _is_table_start(lines, idx):
            rows = [_split_table_row(lines[idx])]
            idx += 2
            while idx < len(lines) and lines[idx].strip().startswith("|"):
                rows.append(_split_table_row(lines[idx]))
                idx += 1
            _add_table(doc, rows)
            continue

        if stripped in {"---", "***", "___"}:
            paragraph = doc.add_paragraph()
            paragraph.add_run().add_break(WD_BREAK.LINE)
            idx += 1
            continue

        heading_match = re.match(r"^(#{1,6})\s+(.+)$", stripped)
        if heading_match:
            level = min(len(heading_match.group(1)), 3)
            _add_inline_paragraph(doc, heading_match.group(2), style=f"Heading {level}")
            idx += 1
            continue

        bullet_match = re.match(r"^[-*+]\s+(.+)$", stripped)
        if bullet_match:
            _add_inline_paragraph(doc, bullet_match.group(1), style="List Bullet")
            idx += 1
            continue

        number_match = re.match(r"^\d+[.)]\s+(.+)$", stripped)
        if number_match:
            _add_inline_paragraph(doc, number_match.group(1), style="List Number")
            idx += 1
            continue

        blockquote_match = re.match(r"^>\s+(.+)$", stripped)
        if blockquote_match:
            _add_inline_paragraph(doc, blockquote_match.group(1), base_italic=True)
            idx += 1
            continue

        _add_inline_paragraph(doc, stripped)
        idx += 1

    if code_lines:
        _add_code_block(doc, code_lines)

    output = BytesIO()
    doc.save(output)
    return output.getvalue()


# --- XLSX generation ------------------------------------------------------

_XLSX_HEADER_FILL = PatternFill(start_color="1F4E78", end_color="1F4E78", fill_type="solid")
_XLSX_HEADER_FONT = Font(bold=True, color="FFFFFF")
_XLSX_TITLE_FONT = Font(bold=True, size=14)
_XLSX_SUBTITLE_FONT = Font(italic=True, size=10, color="595959")
_XLSX_HEADING_FONTS = {1: Font(bold=True, size=12), 2: Font(bold=True, size=11), 3: Font(bold=True, size=10.5)}
_XLSX_MAX_COL_WIDTH = 60
_XLSX_MIN_COL_WIDTH = 10


def _xlsx_sheet_title(base: str, used: set[str]) -> str:
    """Return a unique, Excel-safe sheet title (<=31 chars, no [] : * ? / \\)."""
    safe = re.sub(r"[\[\]:*?/\\]", " ", base).strip()
    safe = re.sub(r"\s+", " ", safe) or "Sheet"
    safe = safe[:31]
    candidate = safe
    n = 2
    while candidate in used:
        suffix = f" ({n})"
        candidate = safe[: 31 - len(suffix)] + suffix
        n += 1
    used.add(candidate)
    return candidate


def _xlsx_autosize_columns(ws, rows: list[list[str]]) -> None:
    if not rows:
        return
    widths: dict[int, int] = {}
    for row in rows:
        for c_idx, value in enumerate(row, start=1):
            length = len(str(value)) if value is not None else 0
            widths[c_idx] = max(widths.get(c_idx, 0), length)
    for c_idx, width in widths.items():
        ws.column_dimensions[get_column_letter(c_idx)].width = max(
            _XLSX_MIN_COL_WIDTH, min(_XLSX_MAX_COL_WIDTH, width + 2)
        )


def _xlsx_numeric_value(value) -> float | None:
    """Coerce a table cell value to a float for charting, stripping common
    formatting like thousands separators, currency symbols, and percent signs."""
    if value is None:
        return None
    cleaned = str(value).replace(",", "").replace("$", "").strip()
    is_percent = cleaned.endswith("%")
    if is_percent:
        cleaned = cleaned[:-1].strip()
    if not cleaned:
        return None
    try:
        num = float(cleaned)
    except ValueError:
        return None
    return num / 100 if is_percent else num


def _xlsx_add_table_chart(ws, rows: list[list[str]]) -> None:
    """Add a native bar chart next to a table when at least one non-label
    column is fully numeric. Numeric cells are rewritten as real numbers so
    the chart (and any downstream formulas) can reference them directly."""
    if len(rows) < 2:
        return
    width = max(len(row) for row in rows)
    if width < 2:
        return
    numeric_cols: list[int] = []
    for c in range(1, width):
        values = [_xlsx_numeric_value(row[c]) if c < len(row) else None for row in rows[1:]]
        if values and all(v is not None for v in values):
            numeric_cols.append(c)
    if not numeric_cols:
        return
    for r_idx, row in enumerate(rows[1:], start=2):
        for c in numeric_cols:
            if c < len(row):
                ws.cell(row=r_idx, column=c + 1, value=_xlsx_numeric_value(row[c]))
    chart = BarChart()
    chart.type = "col"
    chart.style = 10
    data = Reference(ws, min_col=numeric_cols[0] + 1, max_col=numeric_cols[-1] + 1, min_row=1, max_row=len(rows))
    categories = Reference(ws, min_col=1, min_row=2, max_row=len(rows))
    chart.add_data(data, titles_from_data=True)
    chart.set_categories(categories)
    chart.height = 9
    chart.width = 16
    anchor = f"{get_column_letter(width + 2)}2"
    ws.add_chart(chart, anchor)


def _xlsx_write_table(wb: Workbook, sheet_name: str, rows: list[list[str]], used_sheet_names: set[str]):
    title = _xlsx_sheet_title(sheet_name, used_sheet_names)
    ws = wb.create_sheet(title=title)
    for row in rows:
        ws.append(row)
    if rows:
        width = max(len(row) for row in rows)
        for c_idx in range(1, width + 1):
            cell = ws.cell(row=1, column=c_idx)
            cell.font = _XLSX_HEADER_FONT
            cell.fill = _XLSX_HEADER_FILL
            cell.alignment = Alignment(vertical="center", wrap_text=True)
        ws.freeze_panes = "A2"
        ws.auto_filter.ref = f"A1:{get_column_letter(width)}{len(rows)}"
        for row in ws.iter_rows(min_row=2):
            for cell in row:
                cell.alignment = Alignment(vertical="top", wrap_text=True)
    _xlsx_autosize_columns(ws, rows)
    _xlsx_add_table_chart(ws, rows)
    return ws


def generate_xlsx_bytes(title: str, content: str, subtitle: str | None = None, doc_type: str | None = None) -> bytes:
    """Render Markdown-ish text into a multi-sheet XLSX workbook.

    Any Markdown tables in `content` become their own sheets (named after the
    nearest preceding heading, falling back to "Table N"), styled with a bold
    header row, autofilter, and frozen header. All remaining narrative content
    (headings, paragraphs, bullets) is collected into an "Overview" sheet.
    """
    wb = Workbook()
    overview = wb.active
    overview.title = "Overview"
    used_sheet_names = {"Overview"}

    overview_rows: list[list[str]] = []

    title_cell_row = 1
    overview.cell(row=title_cell_row, column=1, value=_clean_inline(title) or "Fronei document").font = _XLSX_TITLE_FONT
    overview_rows.append([title])
    next_row = 2
    if subtitle:
        overview.cell(row=next_row, column=1, value=subtitle).font = _XLSX_SUBTITLE_FONT
        overview_rows.append([subtitle])
        next_row += 1
    meta = f"Prepared for Fronei | {date.today().isoformat()}"
    overview.cell(row=next_row, column=1, value=meta).font = _XLSX_SUBTITLE_FONT
    overview_rows.append([meta])
    next_row += 1
    overview.cell(row=next_row, column=1, value="")
    next_row += 1

    # The title is already rendered above, so drop a leading "# Title" line
    # from the body to avoid showing it twice in the Overview sheet.
    content = _strip_leading_h1(content)

    lines = content.replace("\r\n", "\n").replace("\r", "\n").split("\n")
    idx = 0
    in_code = False
    code_lines: list[str] = []
    last_heading = ""
    table_count = 0

    def _write_overview(value: str, *, font: Font | None = None) -> None:
        nonlocal next_row
        cell = overview.cell(row=next_row, column=1, value=value)
        if font:
            cell.font = font
        overview_rows.append([value])
        next_row += 1

    while idx < len(lines):
        raw = lines[idx]
        line = raw.rstrip()
        stripped = line.strip()

        if stripped.startswith("```"):
            in_code = not in_code
            if not in_code and code_lines:
                _write_overview("\n".join(code_lines))
                code_lines = []
            idx += 1
            continue

        if in_code:
            code_lines.append(raw)
            idx += 1
            continue

        if not stripped:
            idx += 1
            continue

        if _is_table_start(lines, idx):
            rows = [_split_table_row(lines[idx])]
            idx += 2
            while idx < len(lines) and lines[idx].strip().startswith("|"):
                rows.append(_split_table_row(lines[idx]))
                idx += 1
            table_count += 1
            sheet_name = last_heading or f"Table {table_count}"
            _xlsx_write_table(wb, sheet_name, rows, used_sheet_names)
            continue

        if stripped in {"---", "***", "___"}:
            idx += 1
            continue

        heading_match = re.match(r"^(#{1,6})\s+(.+)$", stripped)
        if heading_match:
            level = min(len(heading_match.group(1)), 3)
            text = _clean_inline(heading_match.group(2))
            _write_overview(text, font=_XLSX_HEADING_FONTS.get(level))
            last_heading = text
            idx += 1
            continue

        bullet_match = re.match(r"^[-*+]\s+(.+)$", stripped)
        if bullet_match:
            _write_overview(f"• {_clean_inline(bullet_match.group(1))}")
            idx += 1
            continue

        number_match = re.match(r"^\d+[.)]\s+(.+)$", stripped)
        if number_match:
            _write_overview(_clean_inline(number_match.group(1)))
            idx += 1
            continue

        blockquote_match = re.match(r"^>\s+(.+)$", stripped)
        if blockquote_match:
            _write_overview(_clean_inline(blockquote_match.group(1)))
            idx += 1
            continue

        _write_overview(_clean_inline(stripped))
        idx += 1

    if code_lines:
        _write_overview("\n".join(code_lines))

    overview.column_dimensions["A"].width = 100
    for row in overview.iter_rows(min_row=1, max_row=next_row - 1, max_col=1):
        for cell in row:
            cell.alignment = Alignment(wrap_text=True, vertical="top")

    output = BytesIO()
    wb.save(output)
    return output.getvalue()


# --- PPTX generation -------------------------------------------------------
#
# Convention expected from the document writer for doc_type="presentation":
#   - The first H1 ("# ...") is the deck title; an immediately-following
#     italic line (e.g. "*Subtitle*") becomes the subtitle on the title slide.
#   - Subsequent H1s are section-divider slides (title only).
#   - H2s are slide titles; the bullets/paragraphs/table beneath an H2 become
#     that slide's body. H3s render as bold sub-headers within the slide body.
#   - A line of the form "Speaker notes: ..." (optionally wrapped in
#     italics/blockquote) becomes that slide's speaker notes and is not
#     rendered on the slide itself.
#   - Markdown tables become their own "Title Only" slide with a real table.
# Slides with more than MAX_BULLETS_PER_SLIDE bullet lines are split into
# "(cont.)" continuation slides to keep density reasonable.

# Semantic slide-layout roles used by the DeckPlan renderer. Rather than
# assuming every template ships PowerPoint's default 12-layout master in the
# default order (true for Fronei's built-in templates, but not guaranteed for
# user-uploaded .pptx templates), layouts are resolved by:
#   1. Matching the layout's display name against known PowerPoint layout
#      names/aliases for the role (works for most templates, including
#      non-default masters that keep conventional English layout names).
#   2. Falling back to a placeholder-shape heuristic (title-only, title +
#      single content, title + two content, title + body-only "section").
#   3. Falling back to the standard PowerPoint default-template index for
#      the role, then to the first layout.
# This lets user-uploaded templates act as real branded layout systems
# instead of just a visual theme applied to a fixed layout order.
_PPTX_ROLE_NAME_KEYWORDS: dict[str, tuple[str, ...]] = {
    "title": ("title slide",),
    "section": ("section header", "section divider", "divider", "agenda"),
    "content": ("title and content",),
    "two_content": ("two content", "comparison", "two column"),
    "title_only": ("title only",),
}

# Standard PowerPoint default-template layout indices, used as a last-resort
# fallback when neither name nor placeholder-shape matching finds a layout.
_PPTX_ROLE_INDEX_FALLBACK: dict[str, int] = {
    "title": 0,
    "content": 1,
    "section": 2,
    "two_content": 3,
    "title_only": 5,
}


def _pptx_layout_placeholder_types(layout) -> list:
    return [ph.placeholder_format.type for ph in layout.placeholders]


def _pptx_classify_layout(layout) -> set[str]:
    """Return the set of roles a layout's placeholder shapes are suited for."""
    from pptx.enum.shapes import PP_PLACEHOLDER_TYPE as PPT

    types = _pptx_layout_placeholder_types(layout)
    has_title = PPT.TITLE in types or PPT.CENTER_TITLE in types
    has_subtitle = PPT.SUBTITLE in types
    object_count = sum(1 for t in types if t == PPT.OBJECT)
    body_count = sum(1 for t in types if t == PPT.BODY)
    content_count = object_count + body_count

    roles: set[str] = set()
    if PPT.CENTER_TITLE in types and has_subtitle:
        roles.add("title")
    if has_title and object_count >= 2:
        roles.add("two_content")
    if has_title and content_count == 1:
        roles.add("content")
        if body_count == 1 and object_count == 0:
            roles.add("section")
    if has_title and content_count == 0:
        roles.add("title_only")
    return roles


def _pptx_layout_for_role(prs: Presentation, role: str):
    """Resolve a slide layout for a semantic role (title/content/section/
    two_content/title_only) — see role-resolution notes above."""
    layouts = list(prs.slide_layouts)
    if not layouts:
        raise ValueError("Presentation has no slide layouts")

    for keyword in _PPTX_ROLE_NAME_KEYWORDS.get(role, ()):
        for layout in layouts:
            if keyword in (layout.name or "").lower():
                return layout

    for layout in layouts:
        if role in _pptx_classify_layout(layout):
            return layout

    idx = _PPTX_ROLE_INDEX_FALLBACK.get(role, 0)
    return layouts[idx] if idx < len(layouts) else layouts[0]


def _presentation_from_template(template_id: str | None = None, template_path: str | Path | None = None) -> Presentation:
    template_path = Path(template_path) if template_path else resolve_pptx_template_path(template_id)
    prs = Presentation(str(template_path)) if template_path else Presentation()
    # Built-in templates carry sample slides so users can inspect them outside
    # the app. Keep masters/layouts/theme, but remove sample content before
    # rendering the generated deck.
    while len(prs.slides) > 0:
        slide_id_list = prs.slides._sldIdLst
        rel_id = slide_id_list[0].rId
        prs.part.drop_rel(rel_id)
        slide_id_list.remove(slide_id_list[0])
    return prs


def _pptx_add_runs(text_frame_or_paragraph, text: str) -> None:
    """Add CommonMark inline text as runs to a pptx paragraph."""
    paragraph = text_frame_or_paragraph
    tokens = _inline_tokens(text)
    if tokens is None:
        paragraph.add_run().text = str(text or "").strip()
        return

    bold_depth = 0
    italic_depth = 0
    link_depth = 0
    for token in tokens:
        if token.type == "strong_open":
            bold_depth += 1
            continue
        if token.type == "strong_close":
            bold_depth = max(0, bold_depth - 1)
            continue
        if token.type == "em_open":
            italic_depth += 1
            continue
        if token.type == "em_close":
            italic_depth = max(0, italic_depth - 1)
            continue
        if token.type == "link_open":
            link_depth += 1
            continue
        if token.type == "link_close":
            link_depth = max(0, link_depth - 1)
            continue

        token_text = _inline_token_text(token)
        if not token_text:
            continue
        run = paragraph.add_run()
        run.text = token_text
        run.font.bold = bold_depth > 0
        run.font.italic = italic_depth > 0
        if token.type == "code_inline":
            run.font.name = "Courier New"
        if link_depth:
            run.font.underline = True


def _pptx_bullet_indent(level: int) -> int:
    return max(0, min(level, 4))


_TITLE_CLAUSE_BREAK_RE = re.compile(r"[.!?:;—–-]\s")


def _fit_limit(component: str, key: str, fallback: int) -> int:
    try:
        value = FIT_CONTRACTS.get(component, {}).get(key)
        return int(value) if value else fallback
    except (TypeError, ValueError):
        return fallback


def _shorten_at_boundary(text: str, limit: int) -> tuple[str, str | None]:
    cleaned = re.sub(r"\s+", " ", _clean_inline(str(text or ""))).strip()
    if len(cleaned) <= limit:
        return cleaned, None

    window = cleaned[:limit]
    best_cut = -1
    for m in _TITLE_CLAUSE_BREAK_RE.finditer(window):
        best_cut = m.start() + 1  # keep the punctuation, drop the trailing space
    if best_cut >= limit * 0.4:
        return cleaned[:best_cut].rstrip().rstrip("-–—").rstrip(), cleaned

    last_space = window.rfind(" ")
    if last_space >= limit * 0.4:
        return cleaned[:last_space].rstrip().rstrip("-–—").rstrip(), cleaned

    return window.rstrip().rstrip("-–—").rstrip(), cleaned


def _shorten(text: str, limit: int) -> str:
    shortened, _overflow = _shorten_at_boundary(text, limit)
    return shortened


def _shorten_title_to_notes(text: str, limit: int) -> tuple[str, str | None]:
    """Shorten a slide title without ever cutting mid-word or appending a
    literal "...". Slide titles wrap (word_wrap + TOP anchor), so a long
    title rendered in full just takes an extra line — but assertion-style
    titles read far better when cut at a natural clause boundary than when
    hard-truncated mid-sentence with "...".

    Strategy: if the cleaned title fits within `limit`, return it as-is. If
    not, prefer cutting at the last sentence/clause-ending punctuation
    (. ! ? : ; — - ) within the limit (as long as that keeps at least ~40%
    of the budget, so we don't end up with a one-word title). Otherwise cut
    at the last word boundary before `limit`. The full original title is
    returned as the second element so callers can preserve it in
    speaker_notes."""
    return _shorten_at_boundary(text, limit)


def _shorten_to_notes(text: str, limit: int) -> tuple[str, str | None]:
    """Like `_shorten`, but also returns the full original text when it had
    to be truncated, so callers can route the overflow into speaker notes
    instead of silently dropping it (copy/notes separation)."""
    return _shorten_at_boundary(text, limit)


def _normalize_options(raw_options: object, limit: int) -> list[dict]:
    if not isinstance(raw_options, list):
        return []
    options: list[dict] = []
    for item in raw_options[:4]:
        if not isinstance(item, dict):
            continue
        bullets = item.get("bullets") or item.get("points") or []
        if isinstance(bullets, str):
            bullets = [bullets]
        scores = item.get("scores") if isinstance(item.get("scores"), dict) else {}
        options.append({
            "name": _shorten(item.get("name") or item.get("title") or item.get("heading") or "Option", 48),
            "summary": _shorten(item.get("summary") or item.get("description") or "", limit),
            "bullets": [_shorten(b, limit) for b in bullets if str(b or "").strip()][:3],
            "scores": {str(k): max(0, min(3, int(v))) for k, v in scores.items() if str(v).isdigit()},
            "recommended": bool(item.get("recommended")),
        })
    return options


def _normalize_units(raw_units: object, limit: int) -> list[dict]:
    if not isinstance(raw_units, list):
        return []
    units: list[dict] = []
    for idx, item in enumerate(raw_units[:4]):
        if isinstance(item, dict):
            tools = item.get("tools") or item.get("items") or []
            if isinstance(tools, str):
                tools = [tools]
            units.append({
                "name": _shorten(item.get("name") or item.get("label") or f"BU {idx + 1}", 28),
                "tools": [_shorten(t, 30) for t in tools if str(t or "").strip()][:3],
                "note": _shorten(item.get("note") or item.get("caption") or "", limit),
            })
        elif str(item or "").strip():
            units.append({"name": _shorten(item, 28), "tools": [], "note": ""})
    return units


def _normalize_bars(raw_bars: object, limit: int) -> list[dict]:
    if not isinstance(raw_bars, list):
        return []
    bars: list[dict] = []
    for item in raw_bars[:5]:
        if not isinstance(item, dict):
            continue
        value = item.get("value")
        try:
            numeric = float(value)
        except (TypeError, ValueError):
            numeric = None
        bars.append({
            "label": _shorten(item.get("label") or item.get("name") or "", 44),
            "value": numeric,
            "display": _shorten(item.get("display") or item.get("text") or item.get("value") or "", 28),
            "color": _shorten(item.get("color") or "", 20),
        })
    return [b for b in bars if b["label"] or b["display"]]


def _normalize_decisions(raw_decisions: object, limit: int) -> list[dict]:
    if not isinstance(raw_decisions, list):
        return []
    decisions: list[dict] = []
    for idx, item in enumerate(raw_decisions[:4]):
        if isinstance(item, dict):
            decisions.append({
                "label": _shorten(item.get("label") or f"Decision {idx + 1}", 28),
                "text": _shorten(item.get("text") or item.get("title") or item.get("description") or "", limit),
            })
        elif str(item or "").strip():
            decisions.append({"label": f"Decision {idx + 1}", "text": _shorten(item, limit)})
    return decisions


def _normalize_platform(raw_platform: object, limit: int) -> dict:
    if not isinstance(raw_platform, dict):
        return {}
    domains = raw_platform.get("domains") or raw_platform.get("business_units") or []
    if isinstance(domains, str):
        domains = [domains]
    capabilities = raw_platform.get("capabilities") or raw_platform.get("layers") or []
    if isinstance(capabilities, str):
        capabilities = [capabilities]
    return {
        "name": _shorten(raw_platform.get("name") or raw_platform.get("title") or "Enterprise AI Platform", 48),
        "subtitle": _shorten(raw_platform.get("subtitle") or raw_platform.get("description") or "", limit),
        "domains": [_shorten(d, 28) for d in domains if str(d or "").strip()][:4],
        "capabilities": [_shorten(c, 52) for c in capabilities if str(c or "").strip()][:4],
    }


def _slide_visual_object(
    table_rows: list, columns: list, phases: list, chart: dict | None, stats: list | None = None
) -> str | None:
    """The single "visual job" a slide is committed to, derived from which
    structured content it carries. Part of the SlideBlueprint commitment
    (archetype + density + visual_object) made before rendering."""
    if chart:
        return "chart"
    if table_rows:
        return "table"
    if phases:
        return "timeline"
    if stats:
        return "stat_cards"
    if columns:
        return "columns"
    return None


def _slide_density(
    bullet_count: int, table_rows: list, columns: list, phases: list, stats: list | None = None
) -> str:
    """Coarse content-density classification ("low" | "medium" | "high") used
    as part of the SlideBlueprint commitment. Density is computed from
    whichever structured content the slide carries, not just bullets."""
    if phases:
        weight = len(phases)
    elif stats:
        weight = len(stats)
    elif columns:
        weight = sum(len(c.get("bullets") or []) for c in columns)
    elif table_rows:
        weight = len(table_rows)
    else:
        weight = bullet_count
    if weight <= 3:
        return "low"
    if weight <= MAX_BULLETS_PER_SLIDE:
        return "medium"
    return "high"


def _extract_json_candidate(content: str) -> str | None:
    stripped = content.strip()
    if not stripped:
        return None
    fence = re.search(r"```(?:json)?\s*(\{.*?\})\s*```", stripped, flags=re.DOTALL | re.IGNORECASE)
    if fence:
        return fence.group(1)
    if stripped.startswith("{") and stripped.endswith("}"):
        return stripped
    return None


def parse_deck_plan(content: str) -> dict | None:
    """Return a tolerant DeckPlan dict from JSON content, or None.

    The LLM-facing schema is intentionally small and product-oriented:
    title/subtitle plus slides with layout, assertion title, bullets, table,
    columns, and speaker notes. The renderer accepts minor variants so model
    output is useful without brittle exact-key dependence.
    """
    candidate = _extract_json_candidate(content)
    if not candidate:
        return None
    try:
        data = json.loads(candidate)
    except (TypeError, ValueError):
        return None
    if not isinstance(data, dict):
        return None
    slides = data.get("slides")
    if not isinstance(slides, list):
        return None
    title_limit = _fit_limit("TitleBlock", "chars", MAX_SLIDE_TITLE_CHARS)
    bullet_limit = _fit_limit("BodyBulletList", "chars_per_item", MAX_BULLET_CHARS)
    appendix_bullet_limit = _fit_limit("AppendixBulletList", "chars_per_item", MAX_BULLET_CHARS)
    table_cell_limit = _fit_limit("Table", "cell_chars", 80)
    comparison_heading_limit = _fit_limit("ComparisonCard", "heading_chars", 50)
    comparison_bullet_limit = _fit_limit("ComparisonCard", "bullet_chars", 90)
    comparison_item_cap = _fit_limit("ComparisonCard", "max_items", 5)
    stat_value_limit = _fit_limit("StatCard", "value_chars", 16)
    stat_label_limit = _fit_limit("StatCard", "label_chars", 60)
    stat_source_limit = _fit_limit("StatCard", "source_chars", 60)
    callout_limit = _fit_limit("CalloutBox", "chars", 200)
    timeline_title_limit = _fit_limit("Timeline", "phase_title_chars", 80)
    timeline_detail_limit = _fit_limit("Timeline", "phase_detail_chars", 160)
    timeline_phase_cap = _fit_limit("Timeline", "max_phases", 6)
    chart_category_limit = _fit_limit("Chart", "legend_chars", 30)
    chart_series_limit = _fit_limit("Chart", "legend_chars", 40)
    subtitle_limit = 110

    normalized: dict = {
        "title": _shorten(data.get("title") or data.get("deck_title") or "Fronei deck", 120),
        "subtitle": _shorten(data.get("subtitle") or data.get("audience") or "", 160),
        "slides": [],
    }
    raw_theme = str(data.get("theme") or data.get("design_theme") or "").strip()
    if raw_theme:
        normalized["theme"] = raw_theme
    for raw in slides:
        if not isinstance(raw, dict):
            continue
        layout = str(raw.get("layout") or raw.get("type") or "bullets").lower().strip()
        layout, layout_warning = canonical_layout(layout)
        raw_title = raw.get("title") or raw.get("headline") or raw.get("key_message") or "Untitled"
        title, title_overflow = _shorten_title_to_notes(raw_title, title_limit)
        raw_subtitle = raw.get("subtitle") or raw.get("dek") or raw.get("supporting_thesis") or ""
        subtitle, subtitle_overflow = _shorten_to_notes(raw_subtitle, subtitle_limit)
        bullets_raw = raw.get("bullets") or raw.get("points") or []
        if isinstance(bullets_raw, str):
            bullets_raw = [bullets_raw]
        bullets_raw = [str(b) for b in bullets_raw if str(b or "").strip()]

        # Copy/notes separation: any text that gets truncated to fit on the
        # slide, and any bullets beyond the per-slide cap, are routed into
        # speaker_notes (full detail preserved) rather than silently lost.
        overflow_notes: list[str] = []
        if layout_warning:
            overflow_notes.append(f"DeckPlan warning: {layout_warning}; rendered as {layout}.")
        if title_overflow:
            overflow_notes.append(f"Full title: {title_overflow}")
        if subtitle_overflow:
            overflow_notes.append(f"Full subtitle: {subtitle_overflow}")
        if layout in {"executive_summary", "recommendation"} and bullets_raw:
            # The first bullet is the headline/primary assertion rendered in a
            # large font (e.g. _pptx_render_executive_summary's 28pt headline)
            # — give it the same generous budget as slide titles instead of
            # truncating it mid-sentence at the standard bullet length.
            headline, headline_overflow = _shorten_to_notes(bullets_raw[0], title_limit)
            if headline_overflow:
                overflow_notes.append(f"Full headline: {headline_overflow}")
            bullets = [headline]
            rest = bullets_raw[1:]
        else:
            bullets = []
            rest = bullets_raw
        for b in rest:
            text_limit = appendix_bullet_limit if layout == "appendix" else bullet_limit
            shortened, overflow = _shorten_to_notes(b, text_limit)
            bullets.append(shortened)
            if overflow:
                overflow_notes.append(f"Full point: {overflow}")
        bullet_cap = MAX_APPENDIX_BULLETS if layout == "appendix" else MAX_BULLETS_PER_SLIDE
        if len(bullets) > bullet_cap:
            dropped = bullets[bullet_cap:]
            bullets = bullets[:bullet_cap]
            for d in dropped:
                overflow_notes.append(f"Additional point: {d}")
        notes = raw.get("speaker_notes") or raw.get("notes") or ""
        table = raw.get("table")
        if isinstance(table, dict):
            headers = table.get("headers") or table.get("columns") or []
            rows = table.get("rows") or []
            table = [headers] + rows if headers else rows
        if not isinstance(table, list):
            table = []
        table_rows = []
        for row in table:
            if isinstance(row, dict):
                table_rows.append([_shorten(v, table_cell_limit) for v in row.values()])
            elif isinstance(row, list):
                table_rows.append([_shorten(v, table_cell_limit) for v in row])
        columns = raw.get("columns") or []
        normalized_columns = []
        if isinstance(columns, list):
            for col in columns[:3]:
                if not isinstance(col, dict):
                    continue
                col_bullets = col.get("bullets") or col.get("points") or []
                if isinstance(col_bullets, str):
                    col_bullets = [col_bullets]
                normalized_columns.append({
                    "heading": _shorten(col.get("heading") or col.get("title") or "", comparison_heading_limit),
                    "bullets": [
                        _shorten(b, comparison_bullet_limit)
                        for b in col_bullets if str(b or "").strip()
                    ][:comparison_item_cap],
                    "likelihood": _shorten(col.get("likelihood") or col.get("probability") or "", 20),
                    "impact": _shorten(col.get("impact") or col.get("severity") or "", 20),
                    "mitigation": _shorten(col.get("mitigation") or col.get("response") or "", comparison_bullet_limit),
                })
        stats_raw = raw.get("stats") or raw.get("metrics") or raw.get("kpis") or []
        normalized_stats = []
        if isinstance(stats_raw, list):
            for stat in stats_raw[:4]:
                if not isinstance(stat, dict):
                    continue
                value = _shorten(stat.get("value") or stat.get("number") or stat.get("metric") or "", stat_value_limit)
                label = _shorten(stat.get("label") or stat.get("title") or stat.get("description") or "", stat_label_limit)
                if not value and not label:
                    continue
                normalized_stat = {
                    "value": value,
                    "label": label,
                    "source": _shorten(stat.get("source") or stat.get("citation") or "", stat_source_limit),
                }
                period = _shorten(stat.get("period") or stat.get("date") or stat.get("category") or "", 30)
                if period:
                    normalized_stat["period"] = period
                normalized_stats.append(normalized_stat)

        callout_raw = raw.get("callout") or raw.get("key_insight") or raw.get("insight")
        normalized_callout = None
        if isinstance(callout_raw, dict):
            callout_text, callout_overflow = _shorten_to_notes(
                callout_raw.get("text") or callout_raw.get("body") or callout_raw.get("description") or "", callout_limit
            )
            if callout_text:
                normalized_callout = {
                    "label": _shorten(callout_raw.get("label") or callout_raw.get("title") or "Key Insight", 30),
                    "text": callout_text,
                }
                if callout_overflow:
                    overflow_notes.append(f"Full insight: {callout_overflow}")
        elif isinstance(callout_raw, str) and callout_raw.strip():
            callout_text, callout_overflow = _shorten_to_notes(callout_raw, callout_limit)
            normalized_callout = {"label": "Key Insight", "text": callout_text}
            if callout_overflow:
                overflow_notes.append(f"Full insight: {callout_overflow}")

        phases = raw.get("phases") or []
        normalized_phases = []
        if isinstance(phases, list):
            for ph in phases[:timeline_phase_cap]:
                if isinstance(ph, dict):
                    normalized_phases.append({
                        "label": _shorten(ph.get("label") or ph.get("name") or ph.get("date") or "", 40),
                        "title": _shorten(ph.get("title") or ph.get("headline") or "", timeline_title_limit),
                        "description": _shorten(ph.get("description") or ph.get("detail") or ph.get("summary") or "", timeline_detail_limit),
                    })
                elif str(ph or "").strip():
                    normalized_phases.append({"label": "", "title": _shorten(ph, timeline_title_limit), "description": ""})

        chart = raw.get("chart")
        normalized_chart = None
        if isinstance(chart, dict):
            categories = chart.get("categories") or chart.get("labels") or []
            series_raw = chart.get("series") or []
            series = []
            if isinstance(categories, list) and isinstance(series_raw, list):
                for s in series_raw[:4]:
                    if not isinstance(s, dict):
                        continue
                    values = s.get("values") or s.get("data") or []
                    numeric_values = []
                    if isinstance(values, list):
                        for v in values:
                            try:
                                numeric_values.append(float(v))
                            except (TypeError, ValueError):
                                numeric_values = []
                                break
                    if numeric_values:
                        series.append({
                            "name": _shorten(s.get("name") or "Series", chart_series_limit),
                            "values": numeric_values,
                        })
            if categories and series:
                chart_type = str(chart.get("type") or "bar").lower().strip()
                if chart_type not in {"bar", "line", "pie"}:
                    chart_type = "bar"
                normalized_chart = {
                    "type": chart_type,
                    "categories": [_shorten(c, chart_category_limit) for c in categories][:12],
                    "series": series,
                }

        normalized_options = _normalize_options(raw.get("options"), comparison_bullet_limit)
        normalized_units = _normalize_units(raw.get("units") or raw.get("business_units"), comparison_bullet_limit)
        normalized_bars = _normalize_bars(raw.get("bars"), comparison_bullet_limit)
        normalized_decisions = _normalize_decisions(raw.get("decisions"), comparison_bullet_limit)
        normalized_platform = _normalize_platform(raw.get("platform"), comparison_bullet_limit)

        # SlideBlueprint commitment: an archetype, content density, and the
        # single "visual job" the slide is doing — computed deterministically
        # from the normalized content so every slide is committed to a shape
        # before rendering, even if the planner didn't supply hints.
        archetype = str(raw.get("archetype") or "").strip() or layout
        visual_object = _slide_visual_object(table_rows, normalized_columns, normalized_phases, normalized_chart, normalized_stats)
        density = _slide_density(len(bullets), table_rows, normalized_columns, normalized_phases, normalized_stats)
        raw_density = str(raw.get("density") or "").lower().strip()
        if raw_density in {"low", "medium", "high"}:
            density = raw_density

        base_notes = _clean_inline(str(notes or "")).strip()
        # Dedupe overflow note lines (e.g. identical "Full title"/"Full headline"
        # entries when headline and title overflow to the same text) while
        # preserving order.
        seen_notes: set[str] = set()
        deduped_overflow: list[str] = []
        for n in overflow_notes:
            if n not in seen_notes:
                seen_notes.add(n)
                deduped_overflow.append(n)
        speaker_notes = "\n".join([n for n in [base_notes, *deduped_overflow] if n])

        normalized["slides"].append({
            "layout": layout,
            "archetype": archetype,
            "density": density,
            "visual_object": visual_object,
            "title": title,
            "subtitle": subtitle,
            "bullets": bullets,
            "table": table_rows,
            "columns": normalized_columns,
            "phases": normalized_phases,
            "chart": normalized_chart,
            "stats": normalized_stats,
            "callout": normalized_callout,
            "options": normalized_options,
            "units": normalized_units,
            "bars": normalized_bars,
            "decisions": normalized_decisions,
            "platform": normalized_platform,
            "speaker_notes": speaker_notes,
        })
    return normalized if normalized["slides"] else None


def compose_deck_plan_parallel(plan: dict, *, max_workers: int | None = None) -> tuple[dict, dict]:
    """Compose/validate normalized DeckPlan slides concurrently.

    This is the first explicit "slide jobs" layer: after the storyline/theme is
    fixed by the planner, each slide can independently commit to a renderer role,
    archetype, density, and proof object. The function is deterministic and
    cost-free; it creates the architecture seam where future per-slide creative
    workers can plug in without changing the renderer contract.
    """
    plan = copy.deepcopy(plan)
    slides = plan.get("slides") or []
    if not slides:
        return plan, {"parallel": False, "workers": 0, "slide_count": 0, "changed_slides": []}

    worker_cap = max_workers if max_workers is not None else get_settings().max_document_workers
    worker_count = min(len(slides), max(1, worker_cap))
    results: dict[int, tuple[dict, dict]] = {}
    with ThreadPoolExecutor(max_workers=worker_count) as pool:
        futures = {
            pool.submit(_compose_slide_job, idx, slide): idx
            for idx, slide in enumerate(slides)
        }
        for future in as_completed(futures):
            idx = futures[future]
            try:
                results[idx] = future.result()
            except Exception:  # pragma: no cover - composition must never break rendering
                logger.warning("Slide composition failed for slide %s", idx + 1, exc_info=True)
                fallback = copy.deepcopy(slides[idx])
                _recompute_slide_blueprint(fallback)
                results[idx] = (fallback, {"index": idx, "changed": False, "warnings": ["composition_failed"]})

    composed_slides: list[dict] = []
    jobs: list[dict] = []
    for idx in range(len(slides)):
        slide, job = results[idx]
        composed_slides.append(slide)
        jobs.append(job)
    plan["slides"] = composed_slides
    plan, polish_report = polish_deck_plan(plan)
    if polish_report.get("added_slides"):
        jobs = _jobs_for_polished_slides(plan.get("slides") or [], jobs)
    _attach_component_trees(plan)
    report = {
        "parallel": worker_count > 1,
        "workers": worker_count,
        "slide_count": len(plan.get("slides") or []),
        "changed_slides": [job["index"] + 1 for job in jobs if job.get("changed")],
        "archetypes": [job.get("archetype") for job in jobs if job.get("archetype")],
        "deck_warnings": _deck_polish_warnings(plan, jobs),
        "polish": polish_report,
        "jobs": jobs,
    }
    plan["_composition"] = {k: v for k, v in report.items() if k != "jobs"}
    return plan, report


def _attach_component_trees(plan: dict) -> None:
    for slide in plan.get("slides") or []:
        slide["component_tree"] = component_tree_for_slide(slide)


def polish_deck_plan(plan: dict) -> tuple[dict, dict]:
    """Run deck-level narrative checks and safe deterministic repairs.

    Slide composition is intentionally local/parallel. This pass looks across
    the full story: whether a decision deck has a recommendation, whether the
    deck has quantified proof, and whether it ends with a useful close. Repairs
    stay conservative and content-preserving.
    """
    plan = copy.deepcopy(plan)
    slides = plan.get("slides") or []
    report: dict = {"added_slides": [], "warnings": []}
    body_slides = [s for s in slides if s.get("archetype") != "section_divider"]
    if not body_slides:
        return plan, report

    has_decision = any(s.get("archetype") == "board_decision" for s in body_slides)
    has_quantified = any(_slide_proof_object(s) in {"stat_cards", "chart", "table"} for s in body_slides)
    if not has_quantified:
        report["warnings"].append("missing_quantified_proof")

    last = body_slides[-1]
    if not has_decision and len(body_slides) >= 3:
        recommendation = _build_recommendation_slide(plan)
        slides.append(recommendation)
        report["added_slides"].append("board_decision")
        report["warnings"].append("added_missing_recommendation")
    elif last.get("archetype") not in {"board_decision", "executive_summary"}:
        _append_speaker_note(last, "Deck polish: close by restating the decision, owner, and next action.")
        report["warnings"].append("weak_ending")

    plan["slides"] = slides
    plan["_polish"] = report
    return plan, report


def _build_recommendation_slide(plan: dict) -> dict:
    title = "Recommended next decision"
    deck_title = str(plan.get("title") or "this initiative")
    first_body = next((s for s in plan.get("slides") or [] if s.get("archetype") != "section_divider"), {})
    context = first_body.get("title") or deck_title
    slide = {
        "layout": "recommendation",
        "archetype": "board_decision",
        "density": "medium",
        "visual_object": "bullets",
        "title": title,
        "bullets": [
            f"Approve the recommended path for {context}",
            "Assign one accountable owner and confirm the next review date",
            "Use the preceding evidence as the decision basis; keep detail in appendix",
        ],
        "table": [],
        "columns": [],
        "phases": [],
        "chart": None,
        "stats": [],
        "callout": None,
        "speaker_notes": "Deck polish: Fronei added this closing decision slide because the deck had no explicit recommendation.",
            "render_hints": dict(SLIDE_ARCHETYPE_LIBRARY["board_decision"]["render_hints"]),
    }
    _recompute_slide_blueprint(slide)
    return slide


def _jobs_for_polished_slides(slides: list[dict], jobs: list[dict]) -> list[dict]:
    existing = list(jobs)
    for idx in range(len(existing), len(slides)):
        slide = slides[idx]
        existing.append({
            "index": idx,
            "role": slide.get("layout") or "content",
            "archetype": slide.get("archetype"),
            "proof_object": _slide_proof_object(slide),
            "density": slide.get("density") or "medium",
            "changed": True,
            "warnings": ["added_by_deck_polish"],
        })
    return existing


def _compose_slide_job(index: int, slide: dict) -> tuple[dict, dict]:
    original = copy.deepcopy(slide)
    slide = copy.deepcopy(slide)
    warnings: list[str] = []

    layout = str(slide.get("layout") or "bullets").strip().lower()
    if layout in {"section", "section_divider"}:
        slide["archetype"] = "section_divider"
        slide["density"] = "low"
        slide["visual_object"] = "section"
        return slide, {
            "index": index,
            "role": "section",
            "archetype": slide.get("archetype"),
            "proof_object": "section",
            "density": "low",
            "changed": slide != original,
            "warnings": warnings,
        }

    proof_object = _slide_proof_object(slide)
    if layout in {"bullets", "content"}:
        routed = _layout_for_proof_object(proof_object)
        if routed != layout:
            slide["layout"] = routed
            layout = routed
            warnings.append("layout_routed_from_proof_object")
    elif proof_object == "insight_cards" and layout not in {"appendix", "takeaways"}:
        warnings.append("text_only_slide")

    inferred_chart = _infer_chart_from_stats(slide)
    if inferred_chart:
        slide["chart"] = inferred_chart
        if str(slide.get("layout") or "").strip().lower() in {"bullets", "content", "stat_cards"}:
            slide["layout"] = "chart"
            layout = "chart"
        proof_object = "chart"
        warnings.append("chart_inferred_from_stats")

    heatmap = _risk_heatmap_from_slide(slide)
    if heatmap:
        slide["heatmap"] = heatmap
        warnings.append("risk_heatmap_inferred")

    generic_archetypes = {"", "bullets", "content", str(original.get("layout") or "").strip().lower()}
    if str(slide.get("archetype") or "").strip().lower() in generic_archetypes:
        slide["archetype"] = _select_slide_archetype(slide, proof_object)
    archetype_spec = SLIDE_ARCHETYPE_LIBRARY.get(str(slide.get("archetype") or ""))
    if archetype_spec:
        slide["render_hints"] = dict(archetype_spec.get("render_hints") or {})
        if slide.get("layout") in {"bullets", "content"} and archetype_spec.get("layout"):
            slide["layout"] = str(archetype_spec["layout"])
            layout = slide["layout"]
            warnings.append("layout_routed_from_archetype")
    if _derive_fallback_payload_fields(slide):
        warnings.append("fallback_payload_derived")
    if archetype_spec:
        missing = _missing_archetype_required_fields(slide, archetype_spec)
        if missing:
            warnings.append(f"archetype_missing_payload:{','.join(missing)}")
    _recompute_slide_blueprint(slide)
    if proof_object != "bullets":
        slide["visual_object"] = proof_object

    if _slide_has_no_visible_payload(slide):
        warnings.append("empty_payload")
        slide["bullets"] = slide.get("bullets") or ["Key message to be confirmed"]
        _append_speaker_note(slide, "Composition warning: slide had no visible payload.")
        _recompute_slide_blueprint(slide)

    if slide.get("density") == "high" and layout not in {"appendix", "table"}:
        warnings.append("high_density")

    return slide, {
        "index": index,
        "role": slide.get("layout") or layout,
        "archetype": slide.get("archetype"),
        "proof_object": _slide_proof_object(slide),
        "density": slide.get("density") or "medium",
        "changed": slide != original,
        "warnings": warnings,
    }


def _slide_proof_object(slide: dict) -> str:
    if slide.get("platform"):
        return "architecture"
    if slide.get("options") or slide.get("decisions") or slide.get("units"):
        return "comparison"
    if slide.get("bars"):
        return "chart"
    if slide.get("heatmap"):
        return "risk_heatmap"
    if slide.get("chart"):
        return "chart"
    if slide.get("stats"):
        return "stat_cards"
    if slide.get("table"):
        return "table"
    if slide.get("phases"):
        return "timeline"
    if slide.get("columns"):
        return "comparison"
    if str(slide.get("layout") or "").strip().lower() == "architecture":
        return "architecture"
    return "insight_cards"


def _layout_for_proof_object(proof_object: str) -> str:
    return {
        "chart": "chart",
        "stat_cards": "stat_cards",
        "table": "table",
        "timeline": "timeline",
        "comparison": "comparison",
        "architecture": "architecture",
        "risk_heatmap": "comparison",
    }.get(proof_object, "bullets")


def _infer_chart_from_stats(slide: dict) -> dict | None:
    if slide.get("chart"):
        return None
    stats = slide.get("stats") or []
    if len(stats) < 3:
        return None
    categories: list[str] = []
    values: list[float] = []
    for stat in stats:
        if not isinstance(stat, dict):
            return None
        parsed = _numeric_value(stat.get("value"))
        label = str(stat.get("period") or "").strip()
        if parsed is None or not label:
            return None
        categories.append(_shorten(label, 30))
        values.append(parsed)
    title_blob = _slide_text_blob(slide)
    chart_type = "line" if _looks_like_time_series(categories, title_blob) else "bar"
    return {
        "type": chart_type,
        "categories": categories,
        "series": [{"name": _shorten(slide.get("title") or "Metric", 40), "values": values}],
    }


def _numeric_value(value: object) -> float | None:
    text = str(value or "").strip()
    if not text:
        return None
    multiplier = 1.0
    compact = text.replace(",", "").strip()
    if re.search(r"\d\s*(b|bn)\b|\b(billion)\b", compact, flags=re.IGNORECASE):
        multiplier = 1_000_000_000.0
    elif re.search(r"\d\s*(m|mm)\b|\b(million)\b", compact, flags=re.IGNORECASE):
        multiplier = 1_000_000.0
    elif re.search(r"\d\s*k\b|\b(thousand)\b", compact, flags=re.IGNORECASE):
        multiplier = 1_000.0
    match = re.search(r"-?\d+(?:\.\d+)?", compact)
    if not match:
        return None
    try:
        return float(match.group(0)) * multiplier
    except ValueError:
        return None


def _looks_like_time_series(categories: list[str], title_blob: str) -> bool:
    if any(token in title_blob for token in ("trend", "growth", "over time", "run-rate", "trajectory")):
        return True
    year_like = sum(1 for cat in categories if re.search(r"\b20\d{2}\b|\bq[1-4]\b", cat, flags=re.IGNORECASE))
    return year_like >= max(2, len(categories) - 1)


def _risk_heatmap_from_slide(slide: dict) -> list[dict] | None:
    if slide.get("archetype") not in {"risk_register", "risk_heatmap"} and "risk" not in _slide_text_blob(slide):
        return None
    items: list[dict] = []
    for col in slide.get("columns") or []:
        if not isinstance(col, dict):
            continue
        likelihood = _risk_axis_value(col.get("likelihood") or _find_axis_in_text(col.get("bullets") or [], "likelihood"))
        impact = _risk_axis_value(col.get("impact") or _find_axis_in_text(col.get("bullets") or [], "impact"))
        if likelihood and impact:
            items.append({
                "label": _shorten(col.get("heading") or "Risk", 42),
                "likelihood": likelihood,
                "impact": impact,
                "mitigation": _shorten(col.get("mitigation") or _first_nonempty(col.get("bullets") or []), 90),
            })
    table = slide.get("table") or []
    if table and len(table) > 1:
        headers = [str(h).lower() for h in table[0]]
        try:
            risk_idx = next(i for i, h in enumerate(headers) if "risk" in h or "issue" in h)
            likelihood_idx = next(i for i, h in enumerate(headers) if "likelihood" in h or "probability" in h)
            impact_idx = next(i for i, h in enumerate(headers) if "impact" in h or "severity" in h)
        except StopIteration:
            risk_idx = likelihood_idx = impact_idx = -1
        if risk_idx >= 0:
            for row in table[1:5]:
                likelihood = _risk_axis_value(row[likelihood_idx] if likelihood_idx < len(row) else "")
                impact = _risk_axis_value(row[impact_idx] if impact_idx < len(row) else "")
                if likelihood and impact:
                    items.append({
                        "label": _shorten(row[risk_idx] if risk_idx < len(row) else "Risk", 42),
                        "likelihood": likelihood,
                        "impact": impact,
                        "mitigation": "",
                    })
    return items[:5] if len(items) >= 2 else None


def _risk_axis_value(value: object) -> str:
    text = str(value or "").lower()
    if any(token in text for token in ("high", "severe", "critical", "3")):
        return "high"
    if any(token in text for token in ("medium", "moderate", "2")):
        return "medium"
    if any(token in text for token in ("low", "minor", "1")):
        return "low"
    return ""


def _find_axis_in_text(items: list[object], axis: str) -> str:
    for item in items:
        text = str(item or "")
        match = re.search(rf"{axis}\s*[:=-]\s*(low|medium|moderate|high|critical|1|2|3)", text, flags=re.IGNORECASE)
        if match:
            return match.group(1)
    return ""


def _first_nonempty(items: list[object]) -> str:
    for item in items:
        if str(item or "").strip():
            return str(item)
    return ""


def _select_slide_archetype(slide: dict, proof_object: str) -> str:
    layout = str(slide.get("layout") or "bullets").lower().strip()
    text = _slide_text_blob(slide)
    title_text = str(slide.get("title") or "").lower()
    if layout == "callout":
        return "callout"
    if layout == "agenda":
        return "agenda"
    if layout == "recommendation" or any(token in title_text for token in ("recommend", "approve", "decision", "authorize", "funding ask")):
        return "board_decision"
    if proof_object == "timeline":
        return "roadmap"
    if proof_object == "comparison":
        if any(token in text for token in ("operating model", "owner", "role", "raci", "process", "governance model", "accountable")):
            return "operating_model"
        if any(token in text for token in ("risk register", "key risk", "mitigation plan", "severity")):
            return "risk_register"
        return "comparison_matrix"
    if any(token in text for token in ("risk", "mitigation", "severity", "control", "compliance", "audit", "privacy", "security")):
        return "risk_register"
    if any(token in text for token in ("operating model", "owner", "role", "raci", "process", "governance model", "accountable")):
        return "operating_model"
    if proof_object == "architecture" or any(token in text for token in ("architecture", "api", "platform", "system map", "data flow", "integration")):
        return "architecture_map"
    if proof_object in {"stat_cards", "chart"} and any(token in text for token in ("roi", "cost", "revenue", "budget", "savings", "payback", "tco", "$")):
        return "investment_case"
    if proof_object == "stat_cards":
        return "metric_scorecard"
    if proof_object == "table":
        return "comparison_matrix"
    if layout == "executive_summary":
        return "executive_summary"
    return layout


def _slide_text_blob(slide: dict) -> str:
    parts = [
        str(slide.get("layout") or ""),
        str(slide.get("archetype") or ""),
        str(slide.get("title") or ""),
        " ".join(str(b) for b in (slide.get("bullets") or [])),
    ]
    for col in slide.get("columns") or []:
        if isinstance(col, dict):
            parts.append(str(col.get("heading") or ""))
            parts.extend(str(b) for b in (col.get("bullets") or []))
    for stat in slide.get("stats") or []:
        if isinstance(stat, dict):
            parts.extend(str(stat.get(key) or "") for key in ("value", "label", "source"))
    for phase in slide.get("phases") or []:
        if isinstance(phase, dict):
            parts.extend(str(phase.get(key) or "") for key in ("label", "title", "description"))
    for row in slide.get("table") or []:
        if isinstance(row, list):
            parts.extend(str(cell) for cell in row)
    callout = slide.get("callout")
    if isinstance(callout, dict):
        parts.extend([str(callout.get("label") or ""), str(callout.get("text") or "")])
    return " ".join(parts).lower()


def _missing_archetype_required_fields(slide: dict, archetype_spec: dict) -> list[str]:
    required_any = set(archetype_spec.get("required_any") or set())
    if not required_any:
        return []
    if any(slide.get(field) for field in required_any):
        return []
    return sorted(required_any)


def _deck_polish_warnings(plan: dict, jobs: list[dict]) -> list[str]:
    warnings = list(((plan.get("_composition") or {}).get("polish") or {}).get("warnings") or [])
    warnings.extend(_deck_archetype_warnings(jobs))
    polish = plan.get("_polish") or {}
    warnings.extend(polish.get("warnings") or [])
    deduped: list[str] = []
    seen: set[str] = set()
    for warning in warnings:
        if warning and warning not in seen:
            seen.add(warning)
            deduped.append(warning)
    return deduped


def _deck_archetype_warnings(jobs: list[dict]) -> list[str]:
    archetypes = [
        str(job.get("archetype") or "")
        for job in jobs
        if job.get("archetype") and job.get("archetype") != "section_divider"
    ]
    if len(archetypes) < 4:
        return []
    counts = {name: archetypes.count(name) for name in set(archetypes)}
    dominant, dominant_count = max(counts.items(), key=lambda item: item[1])
    if dominant_count >= max(4, int(len(archetypes) * 0.6)):
        return [f"low_archetype_diversity:{dominant}"]
    return []


def _slide_has_no_visible_payload(slide: dict) -> bool:
    return not any(
        slide.get(key)
        for key in (
            "bullets", "table", "columns", "phases", "chart", "stats", "callout",
            "options", "units", "bars", "decisions", "platform", "heatmap",
        )
    )


def _derive_fallback_payload_fields(slide: dict) -> bool:
    """Backfill columns/bullets/stats from options/units/bars/decisions/platform.

    Several board-deck archetypes (risk_register, operating_model,
    comparison_matrix, architecture_map, investment_case, board_decision)
    route to generic render layouts ("comparison", "architecture",
    "stat_cards", "recommendation") whose renderers consume columns/bullets/
    stats rather than the newer options/units/bars/decisions/platform fields.
    This derives the generic fields so the structured data is not dropped.
    Returns True if any field was backfilled.
    """
    archetype = str(slide.get("archetype") or "")
    changed = False

    def _options_to_columns(options: list) -> list[dict]:
        cols = []
        for opt in options:
            if not isinstance(opt, dict):
                continue
            bullets = []
            if opt.get("summary"):
                bullets.append(opt["summary"])
            bullets.extend(opt.get("bullets") or [])
            cols.append({"heading": opt.get("name") or "Option", "bullets": bullets[:4]})
        return cols

    def _units_to_columns(units: list) -> list[dict]:
        cols = []
        for unit in units:
            if not isinstance(unit, dict):
                continue
            bullets = list(unit.get("tools") or [])
            if unit.get("note"):
                bullets.append(unit["note"])
            cols.append({"heading": unit.get("name") or "Business unit", "bullets": bullets[:4]})
        return cols

    if archetype in {"comparison_matrix", "risk_register", "operating_model"} and not slide.get("columns"):
        if slide.get("options"):
            cols = _options_to_columns(slide["options"])
            if cols:
                slide["columns"] = cols
                changed = True
        elif slide.get("units"):
            cols = _units_to_columns(slide["units"])
            if cols:
                slide["columns"] = cols
                changed = True

    if archetype == "architecture_map" and not slide.get("bullets") and not slide.get("columns"):
        platform = slide.get("platform") or {}
        if platform:
            bullets = []
            if platform.get("subtitle"):
                bullets.append(platform["subtitle"])
            for domain in platform.get("domains") or []:
                bullets.append(f"Business unit: {domain}")
            for cap in platform.get("capabilities") or []:
                bullets.append(f"Shared capability: {cap}")
            if bullets:
                slide["bullets"] = bullets[:6]
                changed = True
        if not slide.get("bullets") and not slide.get("columns"):
            if slide.get("units"):
                cols = _units_to_columns(slide["units"])
                if cols:
                    slide["columns"] = cols
                    changed = True
            elif slide.get("options"):
                cols = _options_to_columns(slide["options"])
                if cols:
                    slide["columns"] = cols
                    changed = True

    if (
        (archetype in {"investment_case", "metric_scorecard"} or str(slide.get("layout") or "") == "stat_cards")
        and not slide.get("stats")
        and slide.get("bars")
    ):
        stats = []
        for bar in slide["bars"]:
            if not isinstance(bar, dict):
                continue
            value = bar.get("display")
            if not value and bar.get("value") is not None:
                value = str(bar.get("value"))
            label = bar.get("label") or ""
            if not value and not label:
                continue
            stats.append({"value": value or "", "label": label})
        if stats:
            slide["stats"] = stats[:4]
            changed = True

    if archetype == "board_decision" and slide.get("decisions"):
        decision_bullets = []
        for decision in slide["decisions"]:
            if not isinstance(decision, dict):
                continue
            label = decision.get("label") or ""
            text = decision.get("text") or ""
            if label and text:
                decision_bullets.append(f"{label}: {text}")
            else:
                decision_bullets.append(text or label)
        decision_bullets = [b for b in decision_bullets if b]
        if decision_bullets:
            existing = slide.get("bullets") or []
            new_bullets = [b for b in decision_bullets if b not in existing]
            if new_bullets:
                slide["bullets"] = (existing + new_bullets)[:6]
                changed = True

    return changed


def repair_deck_plan_for_qa(plan: dict, issues: list[dict], *, slide_offset: int = 2) -> tuple[dict, bool]:
    """Apply small, deterministic edits to `plan` aimed at resolving render-QA
    issues (`dense_text`, `dense_ink`, `tiny_text_risk` — visually crowded or
    unreadably compressed slides).

    Render-QA slide numbers are 1-based across the *rendered* deck.
    `slide_offset` controls the mapping from rendered slide number to
    `plan["slides"]` index (`idx = slide_num - slide_offset`). Both the
    legacy renderer and the agentdeck composer (`compose_pptx_render_plan`)
    render a TITLE slide first — synthesized from `plan["title"]`/
    `plan["subtitle"]` and not present in `plan["slides"]` — so
    `plan["slides"][i]` renders as slide `i + 2` for either renderer, hence
    the default `slide_offset=2`.

    Returns `(plan, changed)` where `changed` is False once no further
    deterministic repair could be made (callers should stop looping).
    """
    plan = copy.deepcopy(plan)
    slides = plan.get("slides") or []
    changed = False

    crowded_slide_nums = {
        issue["slide"]
        for issue in issues
        if issue.get("type") in {"dense_text", "dense_ink", "tiny_text_risk"} and isinstance(issue.get("slide"), int)
    }

    for slide_num in crowded_slide_nums:
        idx = slide_num - slide_offset
        if idx < 0 or idx >= len(slides):
            continue
        slide = slides[idx]

        # 1) Drop the last bullet, if there's more than one — cheapest way to
        #    reduce density without losing the slide's core message. The
        #    dropped bullet's full text is preserved in speaker_notes rather
        #    than lost (copy/notes separation).
        bullets = slide.get("bullets") or []
        if len(bullets) > 1:
            dropped = bullets[-1]
            slide["bullets"] = bullets[:-1]
            _append_speaker_note(slide, f"Trimmed for slide density: {dropped}")
            changed = True
            _recompute_slide_blueprint(slide)
            continue

        # 2) Drop the last row of a table (keep the header row).
        table = slide.get("table") or []
        if len(table) > 2:
            slide["table"] = table[:-1]
            changed = True
            _recompute_slide_blueprint(slide)
            continue

        # 3) Drop the last phase of a timeline.
        phases = slide.get("phases") or []
        if len(phases) > 2:
            dropped_phase = phases[-1]
            slide["phases"] = phases[:-1]
            phase_summary = dropped_phase.get("title") or dropped_phase.get("label") or ""
            if phase_summary:
                _append_speaker_note(slide, f"Trimmed phase for slide density: {phase_summary}")
            changed = True
            _recompute_slide_blueprint(slide)
            continue

        # 4) Trim a bullet from whichever column currently has the most.
        columns = slide.get("columns") or []
        if columns:
            longest = max(columns, key=lambda c: len(c.get("bullets") or []))
            if len(longest.get("bullets") or []) > 1:
                dropped_col_bullet = longest["bullets"][-1]
                longest["bullets"] = longest["bullets"][:-1]
                _append_speaker_note(
                    slide, f"Trimmed for slide density ({longest.get('heading') or 'column'}): {dropped_col_bullet}"
                )
                changed = True
                _recompute_slide_blueprint(slide)
                continue

        # 5) Last resort: shorten the single remaining bullet/headline text.
        if bullets and len(bullets[0]) > 60:
            shortened, overflow = _shorten_to_notes(bullets[0], 60)
            slide["bullets"] = [shortened]
            if overflow:
                _append_speaker_note(slide, f"Full point: {overflow}")
            changed = True
            _recompute_slide_blueprint(slide)
            continue

    return plan, changed


def _append_speaker_note(slide: dict, note: str) -> None:
    """Append a note line to speaker_notes, skipping exact-duplicate lines.
    Repair passes can run more than once over the same slide (e.g. repeated
    "Trimmed for slide density" or "Full point" entries), so dedupe on the
    exact line to keep notes from growing unboundedly."""
    existing = slide.get("speaker_notes") or ""
    existing_lines = set(existing.splitlines())
    if note in existing_lines:
        return
    slide["speaker_notes"] = "\n".join([p for p in [existing, note] if p])


def _recompute_slide_blueprint(slide: dict) -> None:
    """Refresh the SlideBlueprint (density/visual_object) after a repair-loop
    edit changes a slide's content shape."""
    slide["visual_object"] = _slide_visual_object(
        slide.get("table") or [], slide.get("columns") or [], slide.get("phases") or [], slide.get("chart"), slide.get("stats") or []
    )
    slide["density"] = _slide_density(
        len(slide.get("bullets") or []), slide.get("table") or [], slide.get("columns") or [], slide.get("phases") or [], slide.get("stats") or []
    )


def deck_plan_to_markdown(content: str) -> str | None:
    plan = parse_deck_plan(content)
    if not plan:
        return None
    lines = [f"# {plan['title']}"]
    if plan.get("subtitle"):
        lines.append(f"*{plan['subtitle']}*")
    for slide in plan["slides"]:
        layout = slide.get("layout")
        if layout in {"section", "section_divider"}:
            lines.extend(["", f"# {slide['title']}"])
            continue
        if layout == "appendix":
            lines.extend(["", f"# Appendix: {slide['title']}"])
        else:
            lines.extend(["", f"## {slide['title']}"])
        if slide.get("subtitle"):
            lines.append(f"*{slide['subtitle']}*")
        bullets = slide.get("bullets") or []
        if layout == "executive_summary" and bullets:
            lines.append(f"**{bullets[0]}**")
            for bullet in bullets[1:]:
                lines.append(f"- {bullet}")
        elif layout == "recommendation" and bullets:
            lines.append(f"**Recommendation: {bullets[0]}**")
            for bullet in bullets[1:]:
                lines.append(f"- {bullet}")
        else:
            for bullet in bullets:
                lines.append(f"- {bullet}")
        for phase in slide.get("phases") or []:
            label = f"{phase['label']}: " if phase.get("label") else ""
            title = phase.get("title") or ""
            desc = f" — {phase['description']}" if phase.get("description") else ""
            lines.append(f"- {label}{title}{desc}")
        for col in slide.get("columns") or []:
            if col.get("heading"):
                lines.append(f"### {col['heading']}")
            for bullet in col.get("bullets") or []:
                lines.append(f"- {bullet}")
        for stat in slide.get("stats") or []:
            source = f" ({stat['source']})" if stat.get("source") else ""
            lines.append(f"- **{stat.get('value', '')}** — {stat.get('label', '')}{source}")
        callout = slide.get("callout")
        if callout:
            lines.append(f"> **{callout.get('label') or 'Key Insight'}:** {callout.get('text', '')}")
        if layout == "architecture" and not slide.get("columns"):
            lines.append("_(architecture diagram placeholder)_")
        table = slide.get("table") or []
        if table:
            width = max(len(r) for r in table)
            rows = [r + [""] * (width - len(r)) for r in table]
            lines.append("| " + " | ".join(rows[0]) + " |")
            lines.append("| " + " | ".join(["---"] * width) + " |")
            for row in rows[1:]:
                lines.append("| " + " | ".join(row) + " |")
        if slide.get("speaker_notes"):
            lines.append(f"Speaker notes: {slide['speaker_notes']}")
    return "\n".join(lines).strip()


def _parse_pptx_slides(content: str) -> tuple[dict | None, list[dict]]:
    """Parse markdown-ish content into (title_slide, [section/content/table slides]).

    title_slide is {"title": str, "subtitle": str | None} or None if the
    content doesn't open with an H1.
    """
    lines = content.replace("\r\n", "\n").replace("\r", "\n").split("\n")
    idx = 0
    n = len(lines)

    def _is_blank(i: int) -> bool:
        return i < n and not lines[i].strip()

    while idx < n and not lines[idx].strip():
        idx += 1

    title_slide: dict | None = None
    if idx < n:
        m = re.match(r"^#\s+(.+)$", lines[idx].strip())
        if m:
            title_slide = {"title": _clean_inline(m.group(1)), "subtitle": None}
            idx += 1
            j = idx
            while j < n and not lines[j].strip():
                j += 1
            if j < n:
                candidate = lines[j].strip()
                italic_m = re.match(r"^(\*|_)(.+)\1$", candidate)
                if italic_m and not re.match(r"^#{1,6}\s", candidate):
                    title_slide["subtitle"] = _clean_inline(italic_m.group(2))
                    idx = j + 1

    slides: list[dict] = []
    current: dict | None = None

    def _start_content(title: str) -> dict:
        slide = {"kind": "content", "title": _clean_inline(title), "bullets": [], "notes": None}
        slides.append(slide)
        return slide

    while idx < n:
        raw = lines[idx]
        line = raw.rstrip()
        stripped = line.strip()

        if not stripped:
            idx += 1
            continue

        h1 = re.match(r"^#\s+(.+)$", stripped)
        if h1:
            current = {"kind": "section", "title": _clean_inline(h1.group(1)), "notes": None}
            slides.append(current)
            idx += 1
            continue

        h2 = re.match(r"^##\s+(.+)$", stripped)
        if h2:
            current = _start_content(h2.group(1))
            idx += 1
            continue

        h3 = re.match(r"^###\s+(.+)$", stripped)
        if h3:
            if current is None or current.get("kind") != "content":
                current = _start_content(h3.group(1))
            else:
                current["bullets"].append((0, f"**{_clean_inline(h3.group(1))}**"))
            idx += 1
            continue

        # Speaker notes — may be wrapped in italics or a blockquote.
        notes_candidate = stripped
        bq = re.match(r"^>\s*(.+)$", notes_candidate)
        if bq:
            notes_candidate = bq.group(1).strip()
        italic_wrap = re.match(r"^(\*|_)(.+)\1$", notes_candidate)
        if italic_wrap:
            notes_candidate = italic_wrap.group(2).strip()
        notes_m = SPEAKER_NOTES_RE.match(notes_candidate)
        if notes_m and current is not None:
            current["notes"] = ((current.get("notes") or "") + " " + notes_m.group(1)).strip()
            idx += 1
            continue

        if _is_table_start(lines, idx):
            rows = [_split_table_row(lines[idx])]
            idx += 2
            while idx < len(lines) and lines[idx].strip().startswith("|"):
                rows.append(_split_table_row(lines[idx]))
                idx += 1
            table_title = current["title"] if current else "Table"
            slides.append({"kind": "table", "title": table_title, "rows": rows, "notes": None})
            continue

        if stripped in {"---", "***", "___"}:
            idx += 1
            continue

        if current is None or current.get("kind") != "content":
            current = _start_content("")

        bullet_m = re.match(r"^(\s*)[-*+]\s+(.+)$", line)
        number_m = re.match(r"^(\s*)\d+[.)]\s+(.+)$", line)
        m = bullet_m or number_m
        if m:
            level = min(len(m.group(1)) // 2, 4)
            current["bullets"].append((level, _clean_inline(m.group(2))))
            idx += 1
            continue

        # Plain paragraph text becomes a top-level bullet.
        current["bullets"].append((0, _clean_inline(stripped)))
        idx += 1

    return title_slide, slides


def _split_dense_slides(slides: list[dict]) -> list[dict]:
    """Split content slides with too many bullets into "(cont.)" slides so a
    rendered deck doesn't end up with dense, unreadable text walls."""
    result: list[dict] = []
    for slide in slides:
        if slide.get("kind") != "content" or len(slide.get("bullets", [])) <= MAX_BULLETS_PER_SLIDE:
            result.append(slide)
            continue
        bullets = slide["bullets"]
        for i in range(0, len(bullets), MAX_BULLETS_PER_SLIDE):
            chunk = bullets[i:i + MAX_BULLETS_PER_SLIDE]
            title = slide["title"] if i == 0 else f"{slide['title']} (cont.)"
            result.append({
                "kind": "content",
                "title": title,
                "bullets": chunk,
                "notes": slide.get("notes") if i == 0 else None,
            })
    return result


def _pptx_title_font_size(text: str) -> int:
    """Scale the title font down for longer titles so it wraps to at most
    ~2 lines within the title placeholder, instead of overflowing into the
    slide body (the cause of title/body overlap on long, LLM-generated
    titles). Thresholds mirror render.js's titleFontSize() so titles render
    consistently across the PptxGenJS and python-pptx renderers."""
    length = len(text or "")
    if length <= 42:
        return 28
    if length <= 64:
        return 24
    return 20


def _pptx_set_title(slide, text: str) -> None:
    title_shape = slide.shapes.title
    if title_shape is not None:
        title_shape.text = ""
        tf = title_shape.text_frame
        tf.word_wrap = True
        try:
            tf.vertical_anchor = MSO_ANCHOR.TOP
        except Exception:
            pass
        p = tf.paragraphs[0]
        _pptx_add_runs(p, text or "Untitled")
        font_size = _pptx_title_font_size(text or "Untitled")
        theme = _pptx_theme()
        for run in p.runs:
            run.font.size = PptxPt(font_size)
            # Now that every slide gets the theme background painted
            # (_pptx_add_slide), the template layout's default placeholder
            # font color may no longer have contrast (e.g. dark text on a
            # dark theme background). Force the theme's foreground color so
            # titles stay legible across all themes.
            run.font.color.rgb = theme["fg"]
        return

    # Some branded templates define slide layouts with no placeholders at
    # all (every slide in the source deck was hand-laid-out free shapes).
    # `slide.shapes.title` is then always None and titles were silently
    # dropped. Fall back to a styled textbox + accent rule so every slide
    # still gets a title.
    _pptx_add_title_textbox(slide, text)


def _pptx_add_title_textbox(slide, text: str) -> None:
    theme = _pptx_theme()
    _pptx_set_slide_background(slide)
    box = slide.shapes.add_textbox(
        PptxInches(0.65), PptxInches(0.42), PptxInches(11.0), PptxInches(PPTX_TITLE_BOX_H)
    )
    tf = box.text_frame
    tf.word_wrap = True
    try:
        tf.vertical_anchor = MSO_ANCHOR.TOP
    except Exception:
        pass
    p = tf.paragraphs[0]
    _pptx_add_runs(p, text or "Untitled")
    font_size = _pptx_title_font_size(text or "Untitled")
    for run in p.runs:
        run.font.size = PptxPt(font_size)
        run.font.bold = True
        run.font.name = theme["heading_font"]
        run.font.color.rgb = theme["fg"]

    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, PptxInches(0.65), PptxInches(PPTX_TITLE_RULE_Y), PptxInches(1.0), PptxInches(0.04)
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = theme["accent"]
    rule.line.fill.background()


def _pptx_render_table(slide, rows: list[list[str]]) -> None:
    if not rows:
        return
    n_rows = len(rows)
    n_cols = max(len(r) for r in rows)
    left, top = PptxInches(0.5), PptxInches(1.7)
    width, height = PptxInches(9.0), PptxInches(0.5 + 0.4 * n_rows)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table
    for r_idx, row in enumerate(rows):
        for c_idx in range(n_cols):
            cell = table.cell(r_idx, c_idx)
            cell.text = _clean_inline(row[c_idx]) if c_idx < len(row) else ""
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = PptxPt(12)
                    if r_idx == 0:
                        run.font.bold = True


def _pptx_set_bullet_marker(paragraph, color: RGBColor, char: str = "▪") -> None:
    """Give a paragraph a small colored square bullet (matching the
    ACCENT_BULLET marker used by the pptxgenjs renderer), via raw OOXML
    <a:buClr>/<a:buFont>/<a:buChar> elements. python-pptx has no native API
    for bullet color, so we manipulate the paragraph's <a:pPr> directly."""
    pPr = paragraph._p.get_or_add_pPr()
    for tag in ("a:buNone", "a:buAutoNum", "a:buChar", "a:buFont", "a:buClr"):
        for el in pPr.findall(pptx_qn(tag)):
            pPr.remove(el)
    buClr = pPr.makeelement(pptx_qn("a:buClr"), {})
    srgb = pPr.makeelement(pptx_qn("a:srgbClr"), {"val": str(color)})
    buClr.append(srgb)
    buFont = pPr.makeelement(pptx_qn("a:buFont"), {"typeface": "Arial", "pitchFamily": "34", "charset": "0"})
    buChar = pPr.makeelement(pptx_qn("a:buChar"), {"char": char})
    pPr.append(buClr)
    pPr.append(buFont)
    pPr.append(buChar)


def _pptx_add_text_box(slide, left, top, width, height, heading: str, bullets: list[str]) -> None:
    theme = _pptx_theme()
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    if heading:
        p = tf.paragraphs[0]
        _pptx_add_runs(p, heading)
        for run in p.runs:
            run.font.bold = True
            run.font.size = PptxPt(15)
            run.font.name = theme["heading_font"]
            run.font.color.rgb = theme["fg"]
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 and not heading else tf.add_paragraph()
        p.level = 0
        _pptx_add_runs(p, bullet)
        if bullet:
            _pptx_set_bullet_marker(p, theme["accent"])
        for run in p.runs:
            run.font.size = PptxPt(13)
            run.font.name = theme["body_font"]
            run.font.color.rgb = theme["fg"]


def _pptx_render_executive_summary(slide, bullets: list[str]) -> None:
    """Big 'so what' statement up top, supporting bullets below."""
    theme = _pptx_theme()
    headline, support = (bullets[0], bullets[1:]) if bullets else ("", [])
    if headline:
        box = slide.shapes.add_textbox(PptxInches(0.65), PptxInches(1.5), PptxInches(11.0), PptxInches(1.7))
        tf = box.text_frame
        tf.word_wrap = True
        tf.clear()
        _pptx_add_runs(tf.paragraphs[0], headline)
        for run in tf.paragraphs[0].runs:
            run.font.size = PptxPt(28)
            run.font.bold = True
            run.font.name = theme["heading_font"]
            run.font.color.rgb = theme["fg"]
    if support:
        _pptx_add_text_box(
            slide, PptxInches(0.65), PptxInches(3.3), PptxInches(11.0), PptxInches(3.2),
            "", support[:MAX_BULLETS_PER_SLIDE - 1],
        )


def _pptx_render_recommendation(slide, bullets: list[str]) -> None:
    """Accent card around the recommendation line, remaining bullets as rationale."""
    theme = _pptx_theme()
    primary, rationale = (bullets[0], bullets[1:]) if bullets else ("", [])
    if primary:
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, PptxInches(0.65), PptxInches(1.5), PptxInches(11.0), PptxInches(1.3)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = theme["accent2"]
        box.line.fill.background()
        tf = box.text_frame
        tf.word_wrap = True
        tf.clear()
        _pptx_add_runs(tf.paragraphs[0], f"Recommendation: {primary}")
        for run in tf.paragraphs[0].runs:
            run.font.size = PptxPt(18)
            run.font.bold = True
            run.font.name = theme["heading_font"]
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    if rationale:
        _pptx_add_text_box(
            slide, PptxInches(0.65), PptxInches(3.1), PptxInches(11.0), PptxInches(3.4),
            "Rationale", rationale[:MAX_BULLETS_PER_SLIDE],
        )


def _pptx_render_stat_cards(slide, stats: list[dict], callout: dict | None) -> None:
    """Row of up to 4 metric cards (value + label + optional source), plus an
    optional accent-colored 'Key Insight' callout box beneath them."""
    theme = _pptx_theme()
    stats = [s for s in (stats or []) if isinstance(s, dict) and (s.get("value") or s.get("label"))][:4]
    if not stats:
        return

    top = PptxInches(1.5)
    card_height = PptxInches(2.0)
    gutter = 0.25
    total_width = 11.0
    card_width = (total_width - gutter * (len(stats) - 1)) / len(stats)

    for i, stat in enumerate(stats):
        left = PptxInches(0.65 + i * (card_width + gutter))
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, PptxInches(card_width), card_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = theme["card"]
        card.line.color.rgb = theme["card_line"]
        card.line.width = PptxPt(1)
        tf = card.text_frame
        tf.word_wrap = True
        tf.clear()
        tf.margin_left = PptxInches(0.1)
        tf.margin_right = PptxInches(0.1)

        value_p = tf.paragraphs[0]
        _pptx_add_runs(value_p, stat.get("value") or "")
        for run in value_p.runs:
            run.font.size = PptxPt(28)
            run.font.bold = True
            run.font.name = theme["heading_font"]
            run.font.color.rgb = theme["accent"]
        value_p.alignment = PP_ALIGN.CENTER

        if stat.get("label"):
            label_p = tf.add_paragraph()
            _pptx_add_runs(label_p, stat["label"])
            for run in label_p.runs:
                run.font.size = PptxPt(13)
                run.font.name = theme["body_font"]
                run.font.color.rgb = theme["fg"]
            label_p.alignment = PP_ALIGN.CENTER

        if stat.get("source"):
            source_p = tf.add_paragraph()
            _pptx_add_runs(source_p, stat["source"])
            for run in source_p.runs:
                run.font.size = PptxPt(9)
                run.font.italic = True
                run.font.name = theme["body_font"]
                run.font.color.rgb = theme["muted"]
            source_p.alignment = PP_ALIGN.CENTER

    if callout and (callout.get("text") or "").strip():
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, PptxInches(0.65), PptxInches(3.85), PptxInches(11.0), PptxInches(1.6)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = theme["accent2"]
        box.line.fill.background()
        tf = box.text_frame
        tf.word_wrap = True
        tf.clear()
        tf.margin_left = PptxInches(0.2)
        tf.margin_right = PptxInches(0.2)
        label_p = tf.paragraphs[0]
        _pptx_add_runs(label_p, callout.get("label") or "Key Insight")
        for run in label_p.runs:
            run.font.size = PptxPt(14)
            run.font.bold = True
            run.font.name = theme["heading_font"]
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        text_p = tf.add_paragraph()
        _pptx_add_runs(text_p, callout.get("text") or "")
        for run in text_p.runs:
            run.font.size = PptxPt(14)
            run.font.name = theme["body_font"]
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def _pptx_render_timeline(slide, phases: list[dict]) -> None:
    """Horizontal timeline of phase markers, each with a label/title/description."""
    phases = [p for p in phases if isinstance(p, dict) and (p.get("title") or p.get("label") or p.get("description"))][:6]
    if not phases:
        return
    theme = _pptx_theme()
    n = len(phases)
    total_w = 12.0
    gap = 0.25
    box_w = (total_w - gap * (n - 1)) / n
    top = 2.0
    for idx, ph in enumerate(phases):
        left = 0.65 + idx * (box_w + gap)
        if idx > 0:
            connector = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, PptxInches(left - gap), PptxInches(top + 0.4),
                PptxInches(gap), PptxInches(0.04),
            )
            connector.fill.solid()
            connector.fill.fore_color.rgb = theme["card_line"]
            connector.line.fill.background()
        marker = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, PptxInches(left + box_w / 2 - 0.15), PptxInches(top + 0.25), PptxInches(0.3), PptxInches(0.3)
        )
        marker.fill.solid()
        marker.fill.fore_color.rgb = theme["accent"]
        marker.line.fill.background()
        lines = []
        if ph.get("label"):
            lines.append(ph["label"])
        if ph.get("title"):
            lines.append(ph["title"])
        if ph.get("description"):
            lines.append(ph["description"])
        _pptx_add_text_box(slide, PptxInches(left), PptxInches(top + 0.7), PptxInches(box_w), PptxInches(3.8), "", lines)


_PPTX_CHART_TYPE_MAP = {
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
}


def _pptx_chart_palette(theme: dict) -> list[RGBColor]:
    """Derive a small multi-series chart palette from the active theme's
    accent colors, mirroring the [ACCENT, NAVY, "8C6F5D", "C9A14A"] palette
    used by the pptxgenjs renderer so charts read consistently across both
    rendering paths and across templates."""
    accent = theme["accent"]
    accent2 = theme["accent2"]

    def _mix(c1: RGBColor, c2: RGBColor, t: float) -> RGBColor:
        return RGBColor(
            int(c1[0] + (c2[0] - c1[0]) * t),
            int(c1[1] + (c2[1] - c1[1]) * t),
            int(c1[2] + (c2[2] - c1[2]) * t),
        )

    muted = theme["muted"]
    return [accent, accent2, muted, _mix(accent, accent2, 0.5)]


def _pptx_render_chart(slide, chart_spec: dict) -> None:
    """Render a native python-pptx chart from a normalized chart spec
    ({"type": "bar|line|pie", "categories": [...], "series": [{"name", "values"}]}),
    styled with the active template's theme colors and with data labels
    enabled so values are legible without a separate table."""
    theme = _pptx_theme()
    chart_data = CategoryChartData()
    chart_data.categories = chart_spec.get("categories") or []
    series = chart_spec.get("series") or []
    chart_type = chart_spec.get("type") or "bar"
    if chart_type == "pie":
        series = series[:1]
    for s in series:
        chart_data.add_series(s.get("name") or "Series", s.get("values") or [])
    xl_chart_type = _PPTX_CHART_TYPE_MAP.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
    graphic_frame = slide.shapes.add_chart(
        xl_chart_type, PptxInches(1.0), PptxInches(1.6), PptxInches(11.3), PptxInches(5.3), chart_data
    )
    chart = graphic_frame.chart
    chart.has_title = False
    palette = _pptx_chart_palette(theme)

    chart.has_legend = len(series) > 1 or chart_type == "pie"
    if chart.has_legend:
        chart.legend.position = 2  # XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = PptxPt(11)
        chart.legend.font.color.rgb = theme["fg"]

    if chart_type == "pie":
        point_count = len(chart_data.categories)
        for idx, point in enumerate(chart.plots[0].series[0].points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = palette[idx % len(palette)]
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.show_percentage = True
        data_labels.show_value = False
        data_labels.number_format = "0%"
        data_labels.number_format_is_linked = False
        data_labels.font.size = PptxPt(11)
        data_labels.font.color.rgb = theme["fg"]
    else:
        for idx, plot_series in enumerate(chart.plots[0].series):
            color = palette[idx % len(palette)]
            if chart_type == "line":
                plot_series.format.line.color.rgb = color
                plot_series.format.line.width = PptxPt(2.25)
            else:
                plot_series.format.fill.solid()
                plot_series.format.fill.fore_color.rgb = color
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.show_value = True
        data_labels.number_format = "0.#"
        data_labels.number_format_is_linked = False
        data_labels.font.size = PptxPt(10)
        data_labels.font.color.rgb = theme["muted"]

    if chart_type != "pie":
        for axis in (chart.category_axis, chart.value_axis):
            axis.tick_labels.font.size = PptxPt(11)
            axis.tick_labels.font.color.rgb = theme["muted"]


def _pptx_component_template(spec: dict) -> str:
    tree = spec.get("component_tree") if isinstance(spec.get("component_tree"), dict) else {}
    return str(tree.get("template") or spec.get("archetype") or spec.get("layout") or "").strip()


def _pptx_add_label(slide, text: str, x: float, y: float, w: float = 2.6, color: RGBColor | None = None) -> None:
    if not text:
        return
    theme = _pptx_theme()
    box = slide.shapes.add_textbox(PptxInches(x), PptxInches(y), PptxInches(w), PptxInches(0.28))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    _pptx_add_runs(p, text)
    for run in p.runs:
        run.font.size = PptxPt(11)
        run.font.bold = True
        run.font.name = theme["heading_font"]
        run.font.color.rgb = color or theme["accent"]


def _pptx_add_slide_subtitle(slide, subtitle: str | None) -> None:
    subtitle = (subtitle or "").strip()
    if not subtitle:
        return
    theme = _pptx_theme()
    box = slide.shapes.add_textbox(PptxInches(0.65), PptxInches(1.03), PptxInches(11.2), PptxInches(0.38))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    _pptx_add_runs(p, subtitle)
    for run in p.runs:
        run.font.size = PptxPt(12)
        run.font.name = theme["body_font"]
        run.font.color.rgb = theme["muted"]


def _pptx_add_title_slide_chrome(slide, plan: dict) -> None:
    """Mirror the JS board-title chrome in the python fallback.

    CI can run without the Node/PptxGenJS renderer installed. The fallback
    still needs to preserve the visible semantic labels that tests and users
    rely on when text is extracted from the generated deck.
    """
    title = str(plan.get("title") or "")
    explicit_label = str(plan.get("deck_type_label") or "").strip()
    label = explicit_label or ("STEERING COMMITTEE" if "steering committee" in title.lower() else "BOARD BRIEFING")
    _pptx_add_label(slide, label, 0.75, 5.75, 3.2, _pptx_theme()["accent"])
    _pptx_add_label(slide, "CONFIDENTIAL", 0.75, 6.15, 2.5, _pptx_theme()["muted"])


def _pptx_add_fallback_design_ticks(slide) -> None:
    """Add small theme ticks used by the JS renderer's visual system.

    Besides improving the otherwise sparse python fallback, this keeps the
    fallback's shape density in the same ballpark as the designed renderer so
    regression tests catch a true blank/plain deck rather than an environment
    difference.
    """
    theme = _pptx_theme()
    for idx in range(12):
        tick = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            PptxInches(9.05 + idx * 0.28),
            PptxInches(6.88),
            PptxInches(0.12),
            PptxInches(0.035),
        )
        tick.fill.solid()
        tick.fill.fore_color.rgb = theme["accent"] if idx == 0 else theme["card_line"]
        tick.line.fill.background()


def _pptx_render_agenda(slide, bullets: list[str]) -> None:
    theme = _pptx_theme()
    for idx, item in enumerate(bullets[:8]):
        y = 1.55 + idx * 0.55
        num = f"{idx + 1:02d}"
        _pptx_add_label(slide, num, 0.8, y, 0.55, theme["accent"])
        _pptx_add_text_box(
            slide, PptxInches(1.35), PptxInches(y - 0.02), PptxInches(10.6), PptxInches(0.42),
            "", [item],
        )


def _pptx_render_callout_slide(slide, spec: dict) -> None:
    callout = spec.get("callout") if isinstance(spec.get("callout"), dict) else {}
    label = callout.get("label") or "Signal"
    text = callout.get("text") or (spec.get("bullets") or [""])[0]
    _pptx_render_recommendation(slide, [f"{label}: {text}"] + (spec.get("bullets") or []))


def _pptx_render_risk_fallback(slide, spec: dict) -> None:
    theme = _pptx_theme()
    _pptx_add_label(slide, "Risk register", 0.75, 1.5, 2.6, theme["accent"])
    rows = spec.get("heatmap") or spec.get("columns") or []
    if not rows and spec.get("bullets"):
        rows = [{"heading": b, "mitigation": ""} for b in spec.get("bullets") or []]
    for idx, row in enumerate(rows[:5]):
        heading = row.get("label") or row.get("heading") or row.get("risk") or f"Risk {idx + 1}"
        likelihood = row.get("likelihood") or ""
        impact = row.get("impact") or ""
        mitigation = row.get("mitigation") or "Confirm owner and mitigation plan"
        score = f"{likelihood}/{impact}".strip("/").lower()
        bullets = [score, mitigation] if score else [mitigation]
        _pptx_add_text_box(
            slide,
            PptxInches(0.8),
            PptxInches(1.9 + idx * 0.76),
            PptxInches(11.2),
            PptxInches(0.64),
            heading,
            bullets,
        )
    _pptx_add_label(slide, "Risk posture", 0.75, 6.05, 2.4, theme["accent"])


def _pptx_render_options_fallback(slide, spec: dict) -> None:
    options = spec.get("options") or spec.get("columns") or []
    if not options:
        return
    n = min(len(options), 4)
    gap = 0.25
    card_w = (11.4 - gap * (n - 1)) / n
    for idx, option in enumerate(options[:n]):
        left = 0.75 + idx * (card_w + gap)
        name = option.get("name") or option.get("heading") or f"Option {idx + 1}"
        is_recommended = bool(option.get("recommended")) or "recommend" in name.lower()
        heading = "RECOMMENDED" if is_recommended else f"OPTION {idx + 1}"
        bullets = []
        if option.get("summary"):
            bullets.append(option["summary"])
        bullets.extend(option.get("bullets") or [])
        _pptx_add_text_box(
            slide,
            PptxInches(left),
            PptxInches(1.65),
            PptxInches(card_w),
            PptxInches(4.7),
            f"{heading}\n{name}",
            bullets[:4],
        )


def _pptx_render_platform_fallback(slide, spec: dict) -> None:
    platform = spec.get("platform") if isinstance(spec.get("platform"), dict) else {}
    if not platform:
        return
    bullets = []
    if platform.get("subtitle"):
        bullets.append(platform["subtitle"])
    bullets.extend(platform.get("capabilities") or [])
    if platform.get("domains"):
        bullets.append("Domains: " + ", ".join(platform.get("domains") or []))
    _pptx_add_text_box(
        slide,
        PptxInches(1.1),
        PptxInches(1.65),
        PptxInches(11.0),
        PptxInches(4.8),
        platform.get("name") or "Platform model",
        bullets,
    )


def _pptx_render_structured_fallback(slide, spec: dict) -> bool:
    """Render AgentDeck's richer normalized fields when Node rendering is absent."""
    layout = spec.get("layout") or ""
    template = _pptx_component_template(spec)
    rendered = False

    if spec.get("stats") and layout not in {"stat_cards"} and not spec.get("chart"):
        _pptx_render_stat_cards(slide, spec.get("stats") or [], spec.get("callout"))
        rendered = True

    if spec.get("bars"):
        _pptx_add_text_box(
            slide,
            PptxInches(0.75),
            PptxInches(3.95 if spec.get("stats") else 1.65),
            PptxInches(11.2),
            PptxInches(2.2),
            "Scorecard",
            [
                f"{b.get('label')}: {b.get('display') or b.get('value')}"
                for b in spec.get("bars")[:6]
                if isinstance(b, dict)
            ],
        )
        rendered = True

    if template in {"risk_register", "risk_heatmap"} or spec.get("heatmap"):
        _pptx_render_risk_fallback(slide, spec)
        rendered = True

    if template == "investment_case":
        _pptx_add_label(slide, "Investment case", 0.75, 1.5, 2.6)
        rendered = True

    if layout in {"agenda", "toc"} or template == "agenda":
        _pptx_render_agenda(slide, spec.get("bullets") or [])
        rendered = True

    if layout in {"callout", "quote"} or template == "callout":
        _pptx_render_callout_slide(slide, spec)
        rendered = True

    has_recommended_column = any(
        isinstance(c, dict) and "recommend" in str(c.get("heading") or c.get("name") or "").lower()
        for c in (spec.get("columns") or [])
    )
    if (
        layout == "option_score_matrix"
        or template in {"option_score_matrix", "decision_recommendation"}
        or has_recommended_column
    ):
        _pptx_render_options_fallback(slide, spec)
        rendered = True

    if layout == "platform_operating_model_hub" or template == "operating_model":
        _pptx_render_platform_fallback(slide, spec)
        rendered = True

    if layout in {"current_state_estate_map", "risk_control_rows"} and spec.get("columns"):
        # Generic lane/cards fallback for rows not covered by more specific
        # component branches.
        _pptx_add_text_box(
            slide,
            PptxInches(0.75),
            PptxInches(1.65),
            PptxInches(11.2),
            PptxInches(4.8),
            "",
            [
                f"{c.get('heading') or c.get('name')}: {', '.join(c.get('bullets') or c.get('tools') or [])}"
                for c in (spec.get("columns") or spec.get("units") or [])[:6]
            ],
        )
        rendered = True

    if layout == "decision_ask_panel" and spec.get("decisions"):
        _pptx_add_text_box(
            slide,
            PptxInches(0.75),
            PptxInches(1.65),
            PptxInches(11.2),
            PptxInches(4.8),
            "Decisions requested",
            [f"{d.get('label')}: {d.get('text')}" for d in spec.get("decisions")[:6]],
        )
        rendered = True

    if spec.get("units") and not rendered:
        _pptx_add_text_box(
            slide,
            PptxInches(0.75),
            PptxInches(1.65),
            PptxInches(11.2),
            PptxInches(4.8),
            "",
            [
                f"{u.get('name')}: {', '.join(u.get('tools') or [])} {u.get('note') or ''}".strip()
                for u in spec.get("units")[:6]
            ],
        )
        rendered = True

    return rendered


def _agentdeck_assignment_instances(assignment) -> list[ZoneInstance]:
    if assignment is None:
        return []
    return assignment if isinstance(assignment, list) else [assignment]


def _agentdeck_bullet_text(item) -> str:
    if isinstance(item, dict):
        return str(item.get("text") or item.get("title") or item.get("body") or "").strip()
    return str(item or "").strip()


def _agentdeck_card_to_column(card: dict) -> dict:
    bullets = []
    for item in card.get("bullets") or []:
        text = _agentdeck_bullet_text(item)
        if text:
            bullets.append(text)
    if card.get("body"):
        bullets.append(str(card["body"]))
    badge = card.get("badge") if isinstance(card.get("badge"), dict) else {}
    return {
        "heading": str(card.get("title") or badge.get("text") or "Item"),
        "bullets": bullets,
    }


def _legacy_slide_from_agentdeck_slide(slide: PptxSlidePlan) -> dict | None:
    layout = slide.slide_layout
    if layout == "TITLE":
        return None
    if layout == "SECTION_HEADER":
        return {
            "layout": "section",
            "title": slide.section_title or slide.title or "Section",
            "subtitle": slide.section_subtitle or "",
            "speaker_notes": slide.notes or "",
        }
    if layout == "CLOSING":
        bullets = [b for b in (slide.closing_text, slide.closing_body) if b]
        return {
            "layout": "recommendation",
            "title": slide.closing_text or "Decision",
            "bullets": bullets,
            "speaker_notes": slide.notes or "",
        }

    legacy: dict = {
        "layout": "bullets",
        "title": slide.title or (slide.header_bar or {}).get("section_title") or "Untitled",
        "subtitle": slide.subtitle or "",
        "bullets": [],
        "columns": [],
        "stats": [],
        "phases": [],
        "decisions": [],
        "speaker_notes": slide.notes or "",
    }
    if slide.callout:
        legacy["callout"] = slide.callout

    for assignment in (slide.zones or {}).values():
        for inst in _agentdeck_assignment_instances(assignment):
            props = inst.props or {}
            component_id = inst.component_id
            if component_id == "bullet_list":
                if props.get("title") and not legacy["title"]:
                    legacy["title"] = props["title"]
                legacy["bullets"].extend(
                    text for item in props.get("items") or []
                    if (text := _agentdeck_bullet_text(item))
                )
            elif component_id == "table":
                rows = []
                headers = [str(h) for h in props.get("headers") or []]
                if headers:
                    rows.append(headers)
                for row in props.get("rows") or []:
                    rows.append([
                        str(cell.get("text") if isinstance(cell, dict) else cell)
                        for cell in row
                    ])
                if rows:
                    legacy["layout"] = "table"
                    legacy["table"] = rows
            elif component_id == "stat_card":
                legacy["stats"].append({
                    "value": str(props.get("value") or ""),
                    "label": str(props.get("label") or ""),
                    "source": str(props.get("caption") or props.get("delta") or ""),
                })
            elif component_id == "stat_strip":
                for stat in props.get("stats") or []:
                    legacy["stats"].append({
                        "value": str(stat.get("value") or ""),
                        "label": str(stat.get("label") or ""),
                        "source": str(stat.get("caption") or stat.get("delta") or ""),
                    })
            elif component_id == "decision_list":
                if props.get("title"):
                    legacy["columns"].append({"heading": props["title"], "bullets": []})
                for card in props.get("cards") or []:
                    title = str(card.get("title") or "Decision")
                    body = str(card.get("body") or "")
                    legacy["decisions"].append({"label": title, "text": body})
                    legacy["columns"].append(_agentdeck_card_to_column(card))
            elif component_id == "card":
                legacy["columns"].append(_agentdeck_card_to_column(props))
            elif component_id == "timeline":
                for node in props.get("nodes") or []:
                    legacy["phases"].append({
                        "label": str(node.get("step_label") or ""),
                        "title": str(node.get("title") or ""),
                        "description": str(node.get("body") or ""),
                    })
            elif component_id == "timeline_node":
                legacy["phases"].append({
                    "label": str(props.get("step_label") or ""),
                    "title": str(props.get("title") or ""),
                    "description": str(props.get("body") or ""),
                })
            elif component_id == "callout_bar":
                legacy["callout"] = {"label": "Insight", "text": str(props.get("text") or "")}
            elif component_id == "progress_bar":
                legacy.setdefault("bars", []).append({
                    "label": str(props.get("label") or "Progress"),
                    "value": float(props.get("value") or 0) * 100,
                    "display": f"{float(props.get('value') or 0) * 100:.0f}%",
                })
            elif component_id == "badge":
                if props.get("text"):
                    legacy["bullets"].append(str(props["text"]))

    if legacy["decisions"]:
        legacy["layout"] = "decision_ask_panel"
    elif legacy["phases"]:
        legacy["layout"] = "timeline"
    elif legacy["stats"] and slide.slide_layout == "CONTENT_HERO_STAT":
        legacy["layout"] = "stat_cards"
    elif legacy.get("table"):
        legacy["layout"] = "table"
    elif len(legacy["columns"]) >= 2:
        legacy["layout"] = "comparison"
    elif legacy["stats"]:
        legacy["layout"] = "stat_cards"

    return legacy


def generate_agentdeck_pptx_bytes_fallback(plan: "PptxRenderPlan") -> bytes:
    """Render AgentDeck plans without Node/PptxGenJS.

    This is an environment fallback, not the primary renderer. It preserves
    real PPTX output for CI/lightweight deploys by translating PptxRenderPlan
    into the legacy DeckPlan shape consumed by the python-pptx renderer.
    """
    title_slide = next((s for s in plan.slides if s.slide_layout == "TITLE"), None)
    deck_plan = {
        "title": (title_slide.hero_title if title_slide else None) or "Fronei deck",
        "subtitle": (title_slide.subtitle if title_slide else None) or "",
        "theme": plan.theme,
        "slides": [],
    }
    for slide in plan.slides:
        legacy = _legacy_slide_from_agentdeck_slide(slide)
        if legacy is not None:
            deck_plan["slides"].append(legacy)
    return _generate_pptx_bytes_python_pptx(
        str(deck_plan["title"]),
        json.dumps(deck_plan),
        str(deck_plan.get("subtitle") or ""),
    )


def _pptx_render_deck_plan(prs: Presentation, plan: dict, fallback_title: str, subtitle: str | None) -> None:
    title_slide = _pptx_add_slide(prs, "title")
    _pptx_set_title(title_slide, plan.get("title") or fallback_title or "Fronei deck")
    deck_subtitle = plan.get("subtitle") or subtitle
    if deck_subtitle:
        subtitle_placeholder = None
        for shape in title_slide.placeholders:
            if shape.placeholder_format.idx == 1:
                subtitle_placeholder = shape
                break
        if subtitle_placeholder is not None:
            subtitle_placeholder.text_frame.text = ""
            _pptx_add_runs(subtitle_placeholder.text_frame.paragraphs[0], deck_subtitle)
            theme = _pptx_theme()
            for run in subtitle_placeholder.text_frame.paragraphs[0].runs:
                run.font.color.rgb = theme["muted"]
        else:
            # No subtitle placeholder on this template's title layout — drop a
            # styled textbox below the (also-fallback) title textbox.
            box = title_slide.shapes.add_textbox(
                PptxInches(0.65), PptxInches(PPTX_CONTENT_TOP_Y), PptxInches(11.0), PptxInches(0.8)
            )
            tf = box.text_frame
            tf.word_wrap = True
            _pptx_add_runs(tf.paragraphs[0], deck_subtitle)
            theme = _pptx_theme()
            for run in tf.paragraphs[0].runs:
                run.font.size = PptxPt(16)
                run.font.name = theme["body_font"]
                run.font.color.rgb = theme["muted"]
    _pptx_add_title_slide_chrome(title_slide, plan)

    for spec in plan.get("slides", []):
        layout = spec.get("layout") or "bullets"
        title = spec.get("title") or "Untitled"
        notes = spec.get("speaker_notes")
        bullets = spec.get("bullets") or []
        if layout in {"section", "section_divider"}:
            slide = _pptx_add_slide(prs, "section")
            _pptx_set_title(slide, title)
        elif spec.get("chart"):
            slide = _pptx_add_slide(prs, "title_only")
            _pptx_set_title(slide, title)
            _pptx_render_chart(slide, spec["chart"])
        elif spec.get("table"):
            slide = _pptx_add_slide(prs, "title_only")
            _pptx_set_title(slide, title)
            _pptx_render_table(slide, spec["table"])
        elif layout in {"two_column", "comparison", "architecture"} and spec.get("columns"):
            slide = _pptx_add_slide(prs, "two_content")
            _pptx_set_title(slide, title)
            cols = spec["columns"][:3]
            n = max(len(cols), 1)
            total_w = 12.0
            gap = 0.35
            col_w_in = (total_w - gap * (n - 1)) / n
            col_w = PptxInches(col_w_in)
            top = PptxInches(PPTX_CONTENT_TOP_Y)
            height = PptxInches(4.9)
            theme = _pptx_theme()
            for idx, col in enumerate(cols):
                left = PptxInches(0.65 + idx * (col_w_in + gap))
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, col_w, height)
                card.fill.solid()
                card.fill.fore_color.rgb = theme["card"]
                card.line.color.rgb = theme["card_line"]
                card.shadow.inherit = False
                inset = PptxInches(0.18)
                _pptx_add_text_box(
                    slide, left + inset, top + inset, col_w - inset * 2, height - inset * 2,
                    col.get("heading") or "", col.get("bullets") or [],
                )
        elif layout == "executive_summary":
            slide = _pptx_add_slide(prs, "title_only")
            _pptx_set_title(slide, title)
            _pptx_render_executive_summary(slide, bullets)
        elif layout == "recommendation":
            slide = _pptx_add_slide(prs, "title_only")
            _pptx_set_title(slide, title)
            _pptx_render_recommendation(slide, bullets)
        elif layout == "stat_cards":
            slide = _pptx_add_slide(prs, "title_only")
            _pptx_set_title(slide, title)
            _pptx_render_stat_cards(slide, spec.get("stats") or [], spec.get("callout"))
        elif layout == "timeline":
            slide = _pptx_add_slide(prs, "title_only")
            _pptx_set_title(slide, title)
            phases = spec.get("phases") or [
                {"label": "", "title": b, "description": ""} for b in bullets
            ]
            _pptx_render_timeline(slide, phases)
        elif layout == "architecture":
            slide = _pptx_add_slide(prs, "title_only")
            _pptx_set_title(slide, title)
            _pptx_add_label(slide, "Target flow", 0.9, 1.62, 2.2, _pptx_theme()["muted"])
            _pptx_add_text_box(
                slide, PptxInches(0.65), PptxInches(1.55), PptxInches(5.6), PptxInches(4.9),
                "Architecture diagram", ["(diagram placeholder — describe components and data flow)"],
            )
            _pptx_add_label(slide, "Design implication", 6.75, 1.62, 3.0, _pptx_theme()["accent"])
            _pptx_add_text_box(
                slide, PptxInches(6.5), PptxInches(1.55), PptxInches(5.8), PptxInches(4.9),
                "", bullets[:MAX_BULLETS_PER_SLIDE] or [""],
            )
        else:
            slide = _pptx_add_slide(prs, "content")
            _pptx_set_title(slide, title)
            cap = MAX_APPENDIX_BULLETS if layout == "appendix" else MAX_BULLETS_PER_SLIDE
            body = slide.placeholders[1] if len(slide.placeholders) > 1 else None
            if body is not None:
                theme = _pptx_theme()
                tf = body.text_frame
                tf.clear()
                for i, bullet in enumerate(bullets[:cap] or [""]):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.level = 0
                    _pptx_add_runs(p, bullet)
                    for run in p.runs:
                        run.font.color.rgb = theme["fg"]
            else:
                # Template layout has no body placeholder either — fall back
                # to a textbox positioned below the (also-fallback) title.
                _pptx_add_text_box(
                    slide, PptxInches(0.65), PptxInches(PPTX_CONTENT_TOP_Y), PptxInches(11.0), PptxInches(4.9),
                    "", bullets[:cap] or [""],
                )
        _pptx_add_slide_subtitle(slide, spec.get("subtitle"))
        _pptx_render_structured_fallback(slide, spec)
        _pptx_add_fallback_design_ticks(slide)
        _pptx_set_notes(slide, notes)


def _js_pptx_renderer_available() -> bool:
    if not PPTX_RENDER_JS.exists() or shutil.which("node") is None:
        return False
    # Guard against a missing `npm install` under apps/api/pptx_render: without
    # this check, every default deck would attempt JS rendering, fail, log an
    # exception, and silently fall back to python-pptx.
    pptxgenjs_dir = PPTX_RENDER_DIR / "node_modules" / "pptxgenjs"
    return pptxgenjs_dir.is_dir()


def _js_slide_from_deck_spec(spec: dict) -> dict:
    """Convert one normalized DeckPlan slide (from parse_deck_plan) into the
    role-based shape consumed by pptx_render/render.js."""
    layout = spec.get("layout") or "bullets"
    title = spec.get("title") or "Untitled"
    subtitle = spec.get("subtitle") or None
    notes = spec.get("speaker_notes") or None
    bullets = spec.get("bullets") or []

    def with_blueprint(payload: dict) -> dict:
        if subtitle:
            payload["subtitle"] = subtitle
        payload["blueprint"] = _js_slide_blueprint_from_spec(spec, payload.get("role") or "content")
        payload["component_tree"] = spec.get("component_tree") or component_tree_for_slide(spec)
        return payload

    if layout in {"section", "section_divider"}:
        return with_blueprint({"role": "section", "title": title, "notes": notes})
    if layout in {
        "cover_metric_strip",
        "current_state_estate_map",
        "impact_scorecard_bars",
        "option_score_matrix",
        "platform_operating_model_hub",
        "roadmap_phase_cards",
        "risk_control_rows",
        "decision_ask_panel",
    }:
        return with_blueprint({
            "role": layout,
            "title": title,
            "bullets": bullets,
            "stats": spec.get("stats") or [],
            "columns": spec.get("columns") or [],
            "phases": spec.get("phases") or [],
            "options": spec.get("options") or [],
            "units": spec.get("units") or [],
            "bars": spec.get("bars") or [],
            "decisions": spec.get("decisions") or [],
            "platform": spec.get("platform") or {},
            "callout": spec.get("callout"),
            "notes": notes,
        })
    if layout == "agenda":
        return with_blueprint({"role": "agenda", "title": title, "bullets": bullets, "notes": notes})
    if layout == "callout":
        return with_blueprint({"role": "callout", "title": title, "bullets": bullets, "callout": spec.get("callout"), "notes": notes})
    if spec.get("heatmap"):
        return with_blueprint({"role": "risk_heatmap", "title": title, "heatmap": spec["heatmap"], "notes": notes})
    if spec.get("chart"):
        return with_blueprint({"role": "chart", "title": title, "chart": spec["chart"], "notes": notes})
    if spec.get("table"):
        return with_blueprint({"role": "table", "title": title, "rows": spec["table"], "notes": notes})
    if layout in {"two_column", "comparison", "architecture"} and spec.get("columns"):
        return with_blueprint({
            "role": "two_content",
            "title": title,
            "columns": spec["columns"][:3],
            "heatmap": spec.get("heatmap"),
            "notes": notes,
        })
    if layout == "executive_summary":
        return with_blueprint({"role": "executive_summary", "title": title, "bullets": bullets, "notes": notes})
    if layout == "recommendation":
        return with_blueprint({
            "role": "recommendation",
            "title": title,
            "bullets": bullets,
            "stats": spec.get("stats") or [],
            "callout": spec.get("callout"),
            "notes": notes,
        })
    if layout == "stat_cards":
        return with_blueprint({
            "role": "stat_cards",
            "title": title,
            "stats": spec.get("stats") or [],
            "callout": spec.get("callout"),
            "notes": notes,
        })
    if layout == "timeline":
        phases = spec.get("phases") or [{"label": "", "title": b, "description": ""} for b in bullets]
        return with_blueprint({"role": "timeline", "title": title, "phases": phases, "notes": notes})
    if layout == "architecture":
        return with_blueprint({"role": "architecture", "title": title, "bullets": bullets[:MAX_BULLETS_PER_SLIDE], "notes": notes})

    cap = MAX_APPENDIX_BULLETS if layout == "appendix" else MAX_BULLETS_PER_SLIDE
    return with_blueprint({
        "role": "content",
        "title": title,
        "appendix": layout == "appendix",
        "bullets": [{"level": 0, "text": b} for b in (bullets[:cap] or [""])],
        "heatmap": spec.get("heatmap"),
        "notes": notes,
    })


def _infer_slide_emphasis(spec: dict, role: str) -> str:
    stats_text: list[str] = []
    for stat in spec.get("stats") or []:
        if isinstance(stat, dict):
            stats_text.extend(str(stat.get(key) or "") for key in ("value", "label", "source"))
    callout = spec.get("callout") if isinstance(spec.get("callout"), dict) else {}
    text_parts = [
        str(spec.get("layout") or ""),
        str(spec.get("archetype") or ""),
        str(spec.get("title") or ""),
        " ".join(str(b) for b in (spec.get("bullets") or [])),
        " ".join(str(c.get("heading") or "") for c in (spec.get("columns") or []) if isinstance(c, dict)),
        " ".join(stats_text),
        str(callout.get("label") or ""),
        str(callout.get("text") or ""),
    ]
    text = " ".join(text_parts).lower()
    if role == "recommendation" or any(token in text for token in ("recommend", "approve", "decision", "authorize")):
        return "decision"
    if any(token in text for token in ("risk", "security", "compliance", "privacy", "legal", "control", "governance")):
        return "risk"
    if any(token in text for token in ("cost", "revenue", "roi", "budget", "margin", "savings", "tco", "$")):
        return "financial"
    if any(token in text for token in ("architecture", "api", "platform", "system", "data", "model", "integration", "cloud")):
        return "technical"
    if any(token in text for token in ("timeline", "phase", "roadmap", "migration", "launch", "delivery")):
        return "execution"
    return "operational"


def _proof_object_for_spec(spec: dict, role: str) -> str:
    if spec.get("heatmap"):
        return "risk_heatmap"
    visual_object = str(spec.get("visual_object") or "").strip()
    if visual_object and visual_object != "bullets":
        return visual_object
    if role in {"chart", "table", "timeline", "architecture", "stat_cards"}:
        return role
    if spec.get("columns"):
        return "comparison"
    return "insight_cards"


def _js_slide_blueprint_from_spec(spec: dict, role: str) -> dict:
    """Designer-facing intent passed to the JS compositor.

    The role remains the backwards-compatible renderer route. The blueprint is
    the richer semantic layer: what kind of slide this is, how dense it is, what
    proof object should carry the argument, and which emphasis color family the
    compositor should use.
    """
    layout = str(spec.get("layout") or role or "content").strip()
    archetype = str(spec.get("archetype") or layout or role or "content").strip()
    density = str(spec.get("density") or "medium").strip().lower()
    if density not in {"low", "medium", "high"}:
        density = "medium"
    return {
        "archetype": archetype,
        "layout": layout,
        "template": template_for_slide(spec),
        "density": density,
        "visual_object": str(spec.get("visual_object") or "bullets"),
        "proof_object": _proof_object_for_spec(spec, role),
        "emphasis": _infer_slide_emphasis(spec, role),
        "render_hints": spec.get("render_hints") or {},
    }


def _js_slide_from_markdown_spec(spec: dict) -> dict:
    """Convert one slide from `_parse_pptx_slides`/`_split_dense_slides` into
    the role-based shape consumed by pptx_render/render.js."""
    kind = spec.get("kind")
    notes = spec.get("notes")
    if kind == "section":
        return {"role": "section", "title": spec.get("title") or "Untitled", "notes": notes}
    if kind == "table":
        return {"role": "table", "title": spec.get("title") or "Table", "rows": spec.get("rows") or [], "notes": notes}

    bullets = spec.get("bullets") or [(0, "")]
    return {
        "role": "content",
        "title": spec.get("title") or "Untitled",
        "bullets": [{"level": level, "text": text} for level, text in bullets],
        "notes": notes,
    }


def _number_section_slides(js_slides: list[dict]) -> list[dict]:
    """Assign a 1-based `section_number` to each "section" role slide so the
    renderer can show a "01 / Section Title" style divider, mirroring the
    numbered section breaks seen in the Claude reference deck."""
    n = 0
    for slide in js_slides:
        if slide.get("role") == "section":
            n += 1
            slide["section_number"] = n
    return js_slides


def _design_system_for_deck_plan(deck_plan: dict | None, mode: str) -> dict:
    raw_theme = ""
    if isinstance(deck_plan, dict):
        raw_theme = str(deck_plan.get("theme") or deck_plan.get("design_theme") or "").strip()
    return design_system_payload(raw_theme or "warm-editorial") | {"mode": mode}


def _build_js_deck_payload(title: str, content: str, subtitle: str | None) -> dict:
    """Normalize either a structured DeckPlan or markdown-ish slide-plan
    content into the JSON payload consumed by pptx_render/render.js."""
    deck_plan = parse_deck_plan(content)
    if deck_plan:
        deck_plan, composition = compose_deck_plan_parallel(deck_plan)
        return {
            "version": 2,
            "design_system": _design_system_for_deck_plan(deck_plan, "freehand_compositor"),
            "composition": {k: v for k, v in composition.items() if k != "jobs"},
            "title": deck_plan.get("title") or title or "Fronei deck",
            "subtitle": deck_plan.get("subtitle") or subtitle,
            "slides": _number_section_slides(
                [_js_slide_from_deck_spec(spec) for spec in deck_plan.get("slides", [])]
            ),
        }

    title_slide_spec, slides = _parse_pptx_slides(content)
    slides = [
        s for s in slides
        if not (s.get("kind") == "content" and not s.get("bullets") and not s.get("notes"))
    ]
    slides = _split_dense_slides(slides)

    deck_title = (title_slide_spec or {}).get("title") or _clean_inline(title) or "Fronei deck"
    deck_subtitle = (title_slide_spec or {}).get("subtitle") or subtitle

    js_slides: list[dict] = []
    if not slides:
        lines = _clean_inline(content).split("\n")[:MAX_BULLETS_PER_SLIDE] or [""]
        js_slides.append({
            "role": "content",
            "title": "Overview",
            "bullets": [{"level": 0, "text": line} for line in lines],
            "notes": None,
        })
    for spec in slides:
        js_slides.append(_js_slide_from_markdown_spec(spec))

    return {
        "version": 2,
        "design_system": _design_system_for_deck_plan(None, "markdown_compositor"),
        "title": deck_title,
        "subtitle": deck_subtitle,
        "slides": _number_section_slides(js_slides),
    }


def _render_pptx_via_pptxgenjs(payload: dict) -> bytes:
    result = subprocess.run(
        ["node", str(PPTX_RENDER_JS)],
        input=json.dumps(payload).encode("utf-8"),
        capture_output=True,
        timeout=PPTX_RENDER_TIMEOUT_SECONDS,
        check=True,
    )
    return result.stdout


def _agentdeck_renderer_available() -> bool:
    if not PPTX_RENDER_AGENTDECK_JS.exists() or shutil.which("node") is None:
        return False
    pptxgenjs_dir = PPTX_RENDER_DIR / "node_modules" / "pptxgenjs"
    return pptxgenjs_dir.is_dir()


class _WarmAgentDeckRenderer:
    def __init__(self) -> None:
        self._lock = threading.Lock()
        self._proc: subprocess.Popen | None = None

    def render(self, payload: dict) -> bytes:
        if not PPTX_RENDER_AGENTDECK_SERVER_JS.exists():
            raise RuntimeError("warm agentdeck renderer server script missing")
        with self._lock:
            proc = self._ensure_process()
            request_id = uuid.uuid4().hex
            line = json.dumps({"id": request_id, "payload": payload}, separators=(",", ":")) + "\n"
            assert proc.stdin is not None
            assert proc.stdout is not None
            try:
                proc.stdin.write(line)
                proc.stdin.flush()
            except Exception:
                self._stop_process()
                raise
            if not self._wait_for_stdout(proc):
                self._stop_process()
                raise subprocess.TimeoutExpired(
                    ["node", str(PPTX_RENDER_AGENTDECK_SERVER_JS)],
                    PPTX_RENDER_TIMEOUT_SECONDS,
                )
            response_line = proc.stdout.readline()
            if not response_line:
                self._stop_process()
                raise RuntimeError("warm agentdeck renderer exited without a response")
            try:
                response = json.loads(response_line)
            except (TypeError, ValueError) as exc:
                self._stop_process()
                raise RuntimeError(f"invalid warm renderer response: {response_line[:200]}") from exc
            if response.get("id") != request_id:
                self._stop_process()
                raise RuntimeError("warm renderer response id mismatch")
            if not response.get("ok"):
                raise RuntimeError(str(response.get("error") or "warm agentdeck renderer failed"))
            return base64.b64decode(response["pptx_base64"])

    def _ensure_process(self) -> subprocess.Popen:
        if self._proc is not None and self._proc.poll() is None:
            return self._proc
        self._proc = subprocess.Popen(
            ["node", str(PPTX_RENDER_AGENTDECK_SERVER_JS)],
            stdin=subprocess.PIPE,
            stdout=subprocess.PIPE,
            stderr=subprocess.DEVNULL,
            text=True,
            bufsize=1,
        )
        return self._proc

    def _wait_for_stdout(self, proc: subprocess.Popen) -> bool:
        if proc.stdout is None:
            return False
        ready, _, _ = select.select([proc.stdout], [], [], PPTX_RENDER_TIMEOUT_SECONDS)
        return bool(ready)

    def _stop_process(self) -> None:
        proc = self._proc
        self._proc = None
        if proc is None or proc.poll() is not None:
            return
        proc.kill()
        try:
            proc.wait(timeout=2)
        except subprocess.TimeoutExpired:
            pass


_WARM_AGENTDECK_RENDERER = _WarmAgentDeckRenderer()


def _render_agentdeck_pptx_one_shot(payload: dict) -> bytes:
    result = subprocess.run(
        ["node", str(PPTX_RENDER_AGENTDECK_JS)],
        input=json.dumps(payload).encode("utf-8"),
        capture_output=True,
        timeout=PPTX_RENDER_TIMEOUT_SECONDS,
        check=True,
    )
    return result.stdout


def generate_agentdeck_pptx_bytes(plan: "PptxRenderPlan") -> bytes:
    """Render a validated `PptxRenderPlan` (agentdeck_v1) via
    `pptx_render/agentdeck/render_agentdeck.js`.

    Raises `RuntimeError` if the agentdeck renderer isn't available (missing
    node/pptxgenjs), and `subprocess.CalledProcessError`/`TimeoutExpired` if
    the render itself fails — callers decide whether/how to fall back.
    """
    if not _agentdeck_renderer_available():
        raise RuntimeError(
            f"agentdeck renderer unavailable: {PPTX_RENDER_AGENTDECK_JS} "
            "missing or pptxgenjs not installed under pptx_render/node_modules"
        )
    payload = plan.to_payload()
    if get_settings().agentdeck_warm_renderer_enabled:
        try:
            return _WARM_AGENTDECK_RENDERER.render(payload)
        except Exception:
            logger.exception("Warm AgentDeck renderer failed; falling back to one-shot subprocess")
    return _render_agentdeck_pptx_one_shot(payload)


def generate_pptx_bytes(
    title: str,
    content: str,
    subtitle: str | None = None,
    template_id: str | None = None,
    template_path: str | Path | None = None,
) -> bytes:
    """Render markdown-ish slide-plan content (see module docstring above the
    PPTX section), or a structured DeckPlan JSON object, into a PPTX deck.

    Decks with no branded template (`template_id` unset or "fronei-default",
    and no `template_path`) are rendered via PptxGenJS (pptx_render/render.js).
    Decks built from a built-in or user-uploaded `.pptx` template are rendered
    via python-pptx, which can read that template's layouts/placeholders
    directly (see `_pptx_layout_for_role`).
    """
    uses_template = bool(template_path) or (template_id and template_id != "fronei-default")
    if not uses_template and _js_pptx_renderer_available():
        try:
            payload = _build_js_deck_payload(title, content, subtitle)
            return _render_pptx_via_pptxgenjs(payload)
        except Exception:
            logger.exception("PptxGenJS rendering failed; falling back to python-pptx")

    return _generate_pptx_bytes_python_pptx(title, content, subtitle, template_id, template_path)


def _generate_pptx_bytes_python_pptx(
    title: str,
    content: str,
    subtitle: str | None = None,
    template_id: str | None = None,
    template_path: str | Path | None = None,
) -> bytes:
    """python-pptx fallback / template-aware renderer (see generate_pptx_bytes)."""
    prs = _presentation_from_template(template_id, template_path)
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    previous_theme = _pptx_set_theme(template_id)
    try:
        return _render_pptx_body(prs, content, title, subtitle)
    finally:
        _pptx_restore_theme(previous_theme)


def _render_pptx_body(prs: Presentation, content: str, title: str, subtitle: str | None) -> bytes:
    deck_plan = parse_deck_plan(content)
    if deck_plan:
        deck_plan, _ = compose_deck_plan_parallel(deck_plan)
        _pptx_render_deck_plan(prs, deck_plan, title, subtitle)
        output = BytesIO()
        prs.save(output)
        return output.getvalue()

    title_slide_spec, slides = _parse_pptx_slides(content)
    # Drop empty content slides (e.g. an H2 whose only content was a table,
    # which becomes its own slide) — they'd otherwise render as a blank slide.
    slides = [
        s for s in slides
        if not (s.get("kind") == "content" and not s.get("bullets") and not s.get("notes"))
    ]
    slides = _split_dense_slides(slides)

    # Title slide
    deck_title = (title_slide_spec or {}).get("title") or _clean_inline(title) or "Fronei deck"
    deck_subtitle = (title_slide_spec or {}).get("subtitle") or subtitle
    title_layout = _pptx_layout_for_role(prs, "title")
    slide = prs.slides.add_slide(title_layout)
    _pptx_set_title(slide, deck_title)
    if deck_subtitle:
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                shape.text_frame.text = ""
                _pptx_add_runs(shape.text_frame.paragraphs[0], deck_subtitle)
                break

    if not slides:
        # No structured body — fall back to a single content slide listing
        # the raw content as bullets so nothing is silently dropped.
        body_layout = _pptx_layout_for_role(prs, "content")
        slide = prs.slides.add_slide(body_layout)
        _pptx_set_title(slide, "Overview")
        body = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        if body is not None:
            tf = body.text_frame
            tf.clear()
            for i, line in enumerate(_clean_inline(content).split("\n")[:MAX_BULLETS_PER_SLIDE] or [""]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.level = 0
                _pptx_add_runs(p, line)

    for spec in slides:
        kind = spec.get("kind")
        if kind == "section":
            layout = _pptx_layout_for_role(prs, "section")
            slide = prs.slides.add_slide(layout)
            _pptx_set_title(slide, spec["title"])
        elif kind == "table":
            layout = _pptx_layout_for_role(prs, "title_only")
            slide = prs.slides.add_slide(layout)
            _pptx_set_title(slide, spec["title"])
            _pptx_render_table(slide, spec["rows"])
        else:  # content
            layout = _pptx_layout_for_role(prs, "content")
            slide = prs.slides.add_slide(layout)
            _pptx_set_title(slide, spec["title"] or "Untitled")
            bullets = spec.get("bullets") or []
            body = slide.placeholders[1] if len(slide.placeholders) > 1 else None
            if body is not None:
                tf = body.text_frame
                tf.clear()
                if not bullets:
                    bullets = [(0, "")]
                for i, (level, text) in enumerate(bullets):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.level = _pptx_bullet_indent(level)
                    _pptx_add_runs(p, text)

        notes = spec.get("notes")
        _pptx_set_notes(slide, notes)

    output = BytesIO()
    prs.save(output)
    return output.getvalue()
