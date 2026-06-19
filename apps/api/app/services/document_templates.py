from __future__ import annotations

import base64
import logging
import re
import secrets
import shutil
from collections import Counter
from datetime import datetime, timezone
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.config import get_settings
from app.db.models import DocumentTemplate
from app.services.brand.brand_profile import brand_profile_from_template_grammar
from app.services.design_systems.brand_generator import (
    design_system_id_for_template,
    write_brand_design_system,
)

logger = logging.getLogger(__name__)


TEMPLATE_DIR = Path(__file__).resolve().parents[1] / "assets" / "pptx_templates"
MAX_TEMPLATE_UPLOAD_BYTES = 25 * 1024 * 1024

BUILTIN_PPTX_TEMPLATES: dict[str, dict[str, str]] = {
    "fronei-default": {
        "id": "fronei-default",
        "name": "AgentDeck v2",
        "description": "Fronei's default design system, with dark and light variants.",
        "design_system": "agentdeck_v1",
    },
    # Legacy built-ins remain resolvable for older saved briefs and tests, but
    # new presentation setup no longer lists or recommends them.
    "warm-editorial": {
        "id": "warm-editorial",
        "name": "Warm editorial",
        "description": "Thought leadership and executive storytelling with a warm, editorial palette.",
        "filename": "warm_editorial.pptx",
    },
    "modern-tech": {
        "id": "modern-tech",
        "name": "Modern tech",
        "description": "Dark, high-contrast deck for AI platform and product strategy narratives.",
        "filename": "modern_tech.pptx",
    },
    "executive-navy": {
        "id": "executive-navy",
        "name": "Executive navy",
        "description": "Formal boardroom-style navy deck for strategy and architecture reviews.",
        "filename": "executive_navy.pptx",
    },
    "data-product-os": {
        "id": "data-product-os",
        "name": "Data product OS",
        "description": "Dark technical deck for data platforms, governance, and engineering reviews.",
        "filename": "data_product_os.pptx",
    },
    "clean-light": {
        "id": "clean-light",
        "name": "Clean light",
        "description": "Bright, minimal consulting-style deck for operating models and POV decks.",
        "filename": "clean_light.pptx",
    },
}

PREMIUM_FREEHAND_TEMPLATE_ID = "fronei-default"

FREEHAND_SLIDE_TYPES = [
    "hero_cover",
    "executive_summary",
    "decision_recommendation",
    "architecture_map",
    "three_card_system",
    "process_steps",
    "governance_grid",
    "principles_grid",
    "financial_exhibit",
    "risk_matrix",
    "takeaways",
]

TEMPLATE_FOLLOWING_SLIDE_TYPES = [
    "cover",
    "section_divider",
    "executive_summary",
    "two_column_comparison",
    "process_steps",
    "architecture_map",
    "data_exhibit",
    "recommendation",
    "appendix",
]


def resolve_pptx_template_path(template_id: str | None) -> Path | None:
    if not template_id or template_id == "fronei-default":
        return None
    template = BUILTIN_PPTX_TEMPLATES.get(template_id)
    if not template:
        return None
    filename = template.get("filename")
    if not filename:
        return None
    path = TEMPLATE_DIR / filename
    return path if path.exists() else None


def _storage_root() -> Path:
    root = Path(get_settings().document_template_storage_dir).expanduser()
    if not root.is_absolute():
        root = (Path.cwd() / root).resolve()
    root.mkdir(parents=True, exist_ok=True)
    return root


def template_path_for_row(template: DocumentTemplate) -> Path:
    return _storage_root() / template.storage_key


def resolve_template_path(db, user_id: str, template_id: str | None) -> Path | None:
    builtin = resolve_pptx_template_path(template_id)
    if builtin:
        return builtin
    if not template_id or db is None:
        return None
    row = (
        db.query(DocumentTemplate)
        .filter(
            DocumentTemplate.user_id == user_id,
            DocumentTemplate.public_id == template_id,
            DocumentTemplate.is_active == True,  # noqa: E712
        )
        .first()
    )
    if not row:
        return None
    path = template_path_for_row(row)
    return path if path.exists() else None


def _shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    try:
        return (shape.text or "").strip()
    except Exception:
        return ""


def _font_names(prs: Presentation) -> list[str]:
    fonts: Counter[str] = Counter()
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    name = getattr(run.font, "name", None)
                    if name:
                        fonts[name] += 1
    return [name for name, _count in fonts.most_common(5)]


def _rgb_hex_from_color(color) -> str | None:
    try:
        rgb = color.rgb
    except Exception:
        return None
    return str(rgb) if rgb is not None else None


def _theme_colors(prs: Presentation) -> list[str]:
    colors: Counter[str] = Counter()
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                fill = shape.fill
                color = _rgb_hex_from_color(fill.fore_color)
                if color:
                    colors[color] += 1
            except Exception:
                pass
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    color = _rgb_hex_from_color(run.font.color)
                    if color:
                        colors[color] += 1
    return [color for color, _count in colors.most_common(8)]


def _infer_slide_role(slide) -> str:
    shapes = list(slide.shapes)
    text_shapes = [shape for shape in shapes if _shape_text(shape)]
    picture_count = sum(1 for shape in shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE)
    table_count = sum(1 for shape in shapes if getattr(shape, "has_table", False))
    chart_count = sum(1 for shape in shapes if getattr(shape, "has_chart", False))
    text_count = len(text_shapes)
    bulletish_count = sum(1 for shape in text_shapes if "\n" in _shape_text(shape))

    if table_count or chart_count:
        return "data_exhibit"
    if picture_count >= 1 and text_count <= 4:
        return "hero_cover"
    if text_count >= 7:
        return "governance_grid"
    if bulletish_count >= 3 or text_count >= 5:
        return "three_card_system"
    if text_count >= 3:
        return "process_steps"
    return "content"


def _layout_placeholder_types(layout) -> list:
    return [ph.placeholder_format.type for ph in layout.placeholders]


def _placeholder_name(value) -> str:
    try:
        return value.name
    except Exception:
        return str(value)


def _classify_layout_roles(layout) -> set[str]:
    from pptx.enum.shapes import PP_PLACEHOLDER_TYPE as PPT

    types = _layout_placeholder_types(layout)
    has_title = PPT.TITLE in types or PPT.CENTER_TITLE in types
    has_subtitle = PPT.SUBTITLE in types
    object_count = sum(1 for t in types if t == PPT.OBJECT)
    body_count = sum(1 for t in types if t == PPT.BODY)
    content_count = object_count + body_count

    roles: set[str] = set()
    if PPT.CENTER_TITLE in types and has_subtitle:
        roles.add("title")
    if has_title and object_count >= 2:
        roles.add("two_content")
    if has_title and content_count == 1:
        roles.add("content")
        if body_count == 1 and object_count == 0:
            roles.add("section")
    if has_title and content_count == 0:
        roles.add("title_only")
    return roles


def _template_archetypes_for_layout(name: str, roles: list[str], placeholder_types: list[str]) -> list[str]:
    text = " ".join([name, *roles, *placeholder_types]).lower()
    archetypes: list[str] = []
    if "title" in roles or "center_title" in text or "subtitle" in text:
        archetypes.append("cover")
    if "section" in roles or "divider" in text or "agenda" in text:
        archetypes.append("section_divider")
    if "two_content" in roles or "comparison" in text or "two column" in text:
        archetypes.extend(["two_column_comparison", "recommendation"])
    if "content" in roles:
        archetypes.extend(["executive_summary", "process_steps", "architecture_map"])
    if "table" in text or "chart" in text or "object" in text:
        archetypes.extend(["data_exhibit", "risk_matrix", "financial_exhibit"])
    if "blank" in text or "title_only" in roles:
        archetypes.extend(["takeaways", "recommendation"])
    return list(dict.fromkeys(archetypes))


def _layout_inventory(prs: Presentation) -> list[dict[str, object]]:
    inventory: list[dict[str, object]] = []
    for index, layout in enumerate(prs.slide_layouts):
        roles = sorted(_classify_layout_roles(layout))
        placeholder_types = [_placeholder_name(ph.placeholder_format.type) for ph in layout.placeholders]
        placeholder_names = [str(getattr(ph, "name", "") or "") for ph in layout.placeholders if str(getattr(ph, "name", "") or "")]
        inventory.append(
            {
                "index": index,
                "name": layout.name or f"Layout {index}",
                "roles": roles,
                "placeholder_count": len(placeholder_types),
                "placeholder_types": placeholder_types,
                "placeholder_names": placeholder_names[:8],
                "supports": _template_archetypes_for_layout(layout.name or "", roles, placeholder_types),
            }
        )
    return inventory


def _slide_types_from_layout_inventory(inventory: list[dict[str, object]]) -> list[str]:
    found: list[str] = []
    for layout in inventory:
        for archetype in layout.get("supports") or []:
            found.append(str(archetype))
    return list(dict.fromkeys(found))


def _preferred_v3_layouts(inventory: list[dict[str, object]], observed_roles: list[str]) -> list[str]:
    archetypes = set(_slide_types_from_layout_inventory(inventory) + list(observed_roles))
    preferred: list[str] = []
    if {"data_exhibit", "financial_exhibit", "risk_matrix"} & archetypes:
        preferred.append("table")
    if {"two_column_comparison", "recommendation"} & archetypes:
        preferred.extend(["decision", "cards"])
    if {"process_steps", "architecture_map"} & archetypes:
        preferred.append("timeline")
    if {"executive_summary", "takeaways", "three_card_system", "governance_grid"} & archetypes:
        preferred.extend(["cards", "bullets"])
    if {"hero_cover"} & archetypes:
        preferred.append("stat")
    if not preferred:
        preferred = ["cards", "bullets", "table", "timeline", "decision"]
    return list(dict.fromkeys(preferred))


# Cache of computed grammars, keyed by (path, mtime, size) so an edited or
# re-uploaded template (same path, new content) is re-inspected, but repeat
# requests for an unchanged template avoid re-parsing the whole .pptx
# (shapes/fonts/colors/role-inference over every slide) on every presentation
# generation request.
_TEMPLATE_GRAMMAR_CACHE: dict[tuple[str, float, int], dict[str, object]] = {}


MAX_LOGO_ASSET_BYTES = 400 * 1024
_EMU_PER_INCH = 914400


# pptxgenjs's addImage() only reliably handles raster formats it can sniff
# from a data URI; vector/legacy formats commonly embedded as PowerPoint
# "logos" (EMF/WMF from copy-paste, SVG, TIFF, WEBP) make it throw, which
# crashes the whole render subprocess (-> "renderer" stage failure, #198).
_SUPPORTED_LOGO_EXTS = {"png", "jpg", "jpeg", "gif", "bmp"}


def _extract_logo_asset(prs: "Presentation") -> dict[str, object] | None:
    """Best-effort extraction of a brand logo image from an uploaded PPTX
    (#185), used to populate `BrandProfile.logo_assets`.

    Scans the title slide (and first content slide as a fallback) for
    picture shapes and picks the smallest-area picture above a minimum size
    -- logos are usually small marks, whereas large pictures are typically
    full-bleed photos/backgrounds. Returns a JSON-serializable dict with the
    image embedded as base64, or `None` if nothing suitable is found.
    """
    candidates: list[tuple[float, dict[str, object]]] = []
    for slide in list(prs.slides)[:2]:
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            try:
                image = shape.image
            except Exception:
                continue
            if not image.blob or len(image.blob) > MAX_LOGO_ASSET_BYTES:
                continue
            # #198: logos pasted into PowerPoint are frequently EMF/WMF
            # (Windows Metafile), SVG, TIFF, or WEBP -- pptxgenjs's
            # addImage() can't sniff/encode those from a data URI and throws,
            # which crashes the whole render subprocess (renderer-stage
            # DocumentGenerationFailure). Only carry through formats
            # pptxgenjs can actually embed.
            if image.ext.lower() not in _SUPPORTED_LOGO_EXTS:
                continue
            width_in = (shape.width or 0) / _EMU_PER_INCH
            height_in = (shape.height or 0) / _EMU_PER_INCH
            if width_in <= 0 or height_in <= 0:
                continue
            area = width_in * height_in
            # Skip near-full-slide images (backgrounds/photos) -- a logo is
            # a small mark, typically well under a quarter of the slide.
            if area > 6.0:
                continue
            candidates.append((
                area,
                {
                    "content_type": f"image/{image.ext}",
                    "data_base64": base64.b64encode(image.blob).decode("ascii"),
                    "width_in": round(width_in, 3),
                    "height_in": round(height_in, 3),
                },
            ))
    if not candidates:
        return None
    # Smallest qualifying picture is the most likely logo mark.
    candidates.sort(key=lambda c: c[0])
    return candidates[0][1]


def template_grammar_for_path(path: Path) -> dict[str, object]:
    """Inspect a real PPTX template and return a compact design grammar.

    This is intentionally deterministic and cheap. It gives the LLM enough
    guidance to write slides that match the selected template's apparent
    structure before the renderer ever sees the DeckPlan. Results are cached
    per (path, mtime, size) to avoid re-parsing the template on every request.
    """
    try:
        stat = path.stat()
        cache_key = (str(path), stat.st_mtime, stat.st_size)
    except OSError:
        cache_key = None

    if cache_key is not None and cache_key in _TEMPLATE_GRAMMAR_CACHE:
        return _TEMPLATE_GRAMMAR_CACHE[cache_key]

    prs = Presentation(str(path))
    layout_names = [layout.name for layout in prs.slide_layouts if layout.name]
    layout_inventory = _layout_inventory(prs)
    slide_roles = [_infer_slide_role(slide) for slide in prs.slides]
    role_counts = Counter(slide_roles)
    observed_roles = [role for role, _count in role_counts.most_common()]
    logger.info(
        "template_grammar_for_path: inferred slide roles for %s: %s (counts=%s)",
        path.name,
        slide_roles,
        dict(role_counts),
    )
    grammar = {
        "mode": "template_following",
        "source": "pptx_template",
        "slide_count": len(prs.slides),
        "layout_names": layout_names[:12],
        "layout_inventory": layout_inventory[:16],
        "observed_slide_roles": observed_roles[:8],
        "available_slide_types": list(dict.fromkeys(observed_roles + _slide_types_from_layout_inventory(layout_inventory) + TEMPLATE_FOLLOWING_SLIDE_TYPES)),
        "preferred_v3_layouts": _preferred_v3_layouts(layout_inventory, observed_roles),
        "fonts": _font_names(prs),
        "colors": _theme_colors(prs),
        "guidance": [
            "Treat the selected PPTX as a design system, not just a file container.",
            "Use the available slide types repeatedly and consistently rather than inventing a new layout for every slide.",
            "Keep content sparse enough for the template placeholders; move nuance into speaker notes.",
        ],
    }
    if cache_key is not None:
        _TEMPLATE_GRAMMAR_CACHE[cache_key] = grammar
    return grammar


def freehand_premium_grammar(brief: dict | None = None) -> dict[str, object]:
    brief = brief or {}
    return {
        "mode": "fronei_premium_freehand",
        "source": "fronei_default_theme",
        "template_id": PREMIUM_FREEHAND_TEMPLATE_ID,
        "visual_direction": "warm editorial enterprise strategy deck",
        "available_slide_types": FREEHAND_SLIDE_TYPES,
        "fonts": ["Georgia", "Segoe UI"],
        "colors": ["F7F1EE", "E04F00", "282421", "FFFFFF"],
        "guidance": [
            "Use an editorial consulting-deck rhythm: one sharp assertion per slide, visual proof object, short support.",
            "Prefer cards, process diagrams, comparison canvases, and exhibit slides over plain bullet lists.",
            "Use speaker notes for detail that would make the slide crowded.",
        ],
        "recommended_for": brief.get("audience") or "senior stakeholders",
    }


def template_grammar_for_selection(db, user_id: str, template_id: str | None, brief: dict | None = None) -> dict[str, object]:
    if not template_id or template_id == PREMIUM_FREEHAND_TEMPLATE_ID:
        return freehand_premium_grammar(brief)
    path = resolve_template_path(db, user_id, template_id)
    if not path:
        grammar = freehand_premium_grammar(brief)
        grammar["template_id"] = template_id or PREMIUM_FREEHAND_TEMPLATE_ID
        grammar["fallback_reason"] = "selected template was unavailable"
        return grammar
    grammar = template_grammar_for_path(path)
    grammar["template_id"] = template_id
    return grammar


def template_design_context(grammar: dict[str, object] | None) -> str:
    if not grammar:
        return ""
    mode = grammar.get("mode") or "fronei_premium_freehand"
    slide_types = ", ".join(str(x) for x in grammar.get("available_slide_types") or []) or "executive_summary, recommendation"
    fonts = ", ".join(str(x) for x in grammar.get("fonts") or []) or "template fonts"
    colors = ", ".join(str(x) for x in grammar.get("colors") or []) or "template palette"
    guidance = "\n".join(f"- {item}" for item in (grammar.get("guidance") or []))
    if mode == "template_following":
        layout_names = ", ".join(str(x) for x in grammar.get("layout_names") or [])
        observed = ", ".join(str(x) for x in grammar.get("observed_slide_roles") or [])
        return f"""TEMPLATE-FIRST PRESENTATION DESIGN BRIEF:
Mode: follow the selected PowerPoint template.
Selected template id: {grammar.get("template_id") or "unknown"}
Observed template slide roles: {observed or "not enough sample slides to infer; use conventional executive deck roles"}.
Available template layout names: {layout_names or "unavailable"}.
Use only these semantic slide types unless the user explicitly asks otherwise: {slide_types}.
Template fonts observed: {fonts}.
Template colors observed: {colors}.
{guidance}

DeckPlan rules for this template:
- Do not produce a Markdown article in slide form. Produce a designed presentation plan.
- Each slide must have one assertion-style title under 12 words and one visual job: compare, explain flow, quantify impact, show timeline, or ask for a decision.
- Prefer `executive_summary`, `recommendation`, `timeline`, `architecture`, `comparison`, `financial_model`, `risk_matrix`, and `table` layouts over generic `bullets`.
- Keep visible slide text sparse: max 3 bullets per slide, max 12 words per bullet, no paragraph bullets.
- Put detail, caveats, and narration in `speaker_notes`, not on the slide."""

    return f"""TEMPLATE-FIRST PRESENTATION DESIGN BRIEF:
Mode: Fronei premium freehand theme.
Visual direction: {grammar.get("visual_direction") or "editorial enterprise strategy deck"}.
Use these semantic slide types: {slide_types}.
Theme fonts: {fonts}.
Theme colors: {colors}.
{guidance}

DeckPlan rules for this theme:
- Think like a human presentation designer: define the story spine, choose a slide archetype for each moment, then write sparse slide copy.
- Do not produce generic bullet-only decks. Every slide needs a visual structure: callout, cards, matrix, roadmap, architecture map, chart, or decision panel.
- Each slide title must be an assertion under 12 words. Avoid centered, multi-line newspaper headlines.
- Visible text limit: max 3 bullets per slide, max 12 words per bullet. Move nuance into `speaker_notes`.
- For business/strategy decks, include an executive summary, evidence/analysis section, decision/recommendation slide, and next steps."""


def _safe_name(name: str) -> str:
    cleaned = re.sub(r"\s+", " ", (name or "").strip())
    return cleaned[:160] or "Presentation template"


def _template_option_from_row(row: DocumentTemplate, *, recommended: bool = False) -> dict[str, object]:
    return {
        "id": row.public_id,
        "name": row.name,
        "description": row.description or f"Uploaded from {row.original_filename or 'PowerPoint template'}",
        "recommended": recommended,
        "user_template": True,
        "design_mode": "template_following",
        # #185: the brand design_system generated from this template's
        # BrandProfile (#181/#184), if any -- lets the picker show a "Brand"
        # badge and #183 resolves this id at generation time.
        "design_system": row.design_system_id or None,
    }


def recommend_template_id(brief: dict | None) -> str:
    return "fronei-default"


def list_document_templates(
    doc_type: str | None = None,
    brief: dict | None = None,
    db=None,
    user_id: str | None = None,
) -> list[dict[str, object]]:
    recommendation = recommend_template_id(brief)
    if doc_type != "presentation":
        base = BUILTIN_PPTX_TEMPLATES["fronei-default"].copy()
        base["recommended"] = True
        base["design_mode"] = "fronei_premium_freehand"
        return [base]

    templates: list[dict[str, object]] = []
    user_rows: list[DocumentTemplate] = []
    if db is not None and user_id:
        user_rows = (
            db.query(DocumentTemplate)
            .filter(
                DocumentTemplate.user_id == user_id,
                DocumentTemplate.doc_type == "presentation",
                DocumentTemplate.is_active == True,  # noqa: E712
            )
            .order_by(DocumentTemplate.updated_at.desc())
            .all()
        )
        templates.extend(_template_option_from_row(row) for row in user_rows)

    for template_id in ("fronei-default",):
        item = BUILTIN_PPTX_TEMPLATES[template_id].copy()
        if template_id != "fronei-default" and not resolve_pptx_template_path(template_id):
            continue
        item["recommended"] = not user_rows and template_id == recommendation
        item["design_mode"] = "fronei_premium_freehand" if template_id == "fronei-default" else "template_following"
        templates.append(item)
    if user_rows and templates:
        templates[0]["recommended"] = True
    elif not any(t.get("recommended") for t in templates) and templates:
        templates[0]["recommended"] = True
    return templates


def store_user_pptx_template(
    db,
    user_id: str,
    *,
    filename: str,
    content_type: str | None,
    data: bytes,
    name: str | None = None,
    description: str | None = None,
) -> DocumentTemplate:
    if not data:
        raise ValueError("Template file is empty.")
    if len(data) > MAX_TEMPLATE_UPLOAD_BYTES:
        raise ValueError("Template file is too large.")
    if not (filename or "").lower().endswith(".pptx"):
        raise ValueError("Only .pptx templates are supported.")

    # Validate before writing permanently.
    prs = Presentation(BytesIO(data))
    logo_asset = None
    try:
        logo_asset = _extract_logo_asset(prs)
    except Exception:
        logger.exception("Failed to extract logo asset from uploaded template %s", filename)

    now = datetime.now(timezone.utc)
    public_id = secrets.token_hex(12)
    user_dir = _storage_root() / user_id
    user_dir.mkdir(parents=True, exist_ok=True)
    storage_key = f"{user_id}/{public_id}.pptx"
    path = _storage_root() / storage_key
    with path.open("wb") as f:
        f.write(data)

    row = DocumentTemplate(
        public_id=public_id,
        user_id=user_id,
        name=_safe_name(name or Path(filename).stem),
        description=(description or "").strip()[:500] or None,
        doc_type="presentation",
        storage_key=storage_key,
        original_filename=filename[:255],
        content_type=(content_type or "")[:120] or None,
        file_size=len(data),
        is_active=True,
        created_at=now,
        updated_at=now,
    )
    db.add(row)
    db.commit()
    db.refresh(row)

    # #184: generate a first-class brand design_system from this template's
    # grammar/colors/fonts (#181) and register it on the row so the document
    # popup can select it (#183 resolves it at generation time). Best-effort:
    # a failure here should not block the template upload itself.
    try:
        grammar = template_grammar_for_path(path)
        brand_profile = brand_profile_from_template_grammar(
            grammar, user_id=user_id, profile_id=f"template:{public_id}"
        )
        brand_profile.source_template_id = public_id
        if logo_asset is not None:
            brand_profile.logo_assets = [logo_asset]
        design_system_id = design_system_id_for_template(user_id, public_id)
        write_brand_design_system(brand_profile, design_system_id=design_system_id)
        row.design_system_id = design_system_id
        row.updated_at = datetime.now(timezone.utc)
        db.commit()
        db.refresh(row)
    except Exception:
        logger.exception("Failed to generate brand design_system for template %s", public_id)

    return row


def rename_user_template(db, user_id: str, template_id: str, name: str) -> DocumentTemplate | None:
    row = (
        db.query(DocumentTemplate)
        .filter(
            DocumentTemplate.user_id == user_id,
            DocumentTemplate.public_id == template_id,
            DocumentTemplate.is_active == True,  # noqa: E712
        )
        .first()
    )
    if not row:
        return None
    safe = _safe_name(name)
    if not safe:
        raise ValueError("Template name cannot be empty.")
    row.name = safe
    row.updated_at = datetime.now(timezone.utc)
    db.commit()
    db.refresh(row)
    return row


def replace_user_pptx_template(
    db,
    user_id: str,
    template_id: str,
    *,
    filename: str,
    content_type: str | None,
    data: bytes,
) -> DocumentTemplate | None:
    """Re-upload the PPTX content for an existing template row, keeping the
    same `public_id`/identity (so any saved references to it stay valid),
    while replacing the stored file and regenerating the brand design_system
    (#184) from the new content's grammar/colors/fonts/logo.
    """
    if not data:
        raise ValueError("Template file is empty.")
    if len(data) > MAX_TEMPLATE_UPLOAD_BYTES:
        raise ValueError("Template file is too large.")
    if not (filename or "").lower().endswith(".pptx"):
        raise ValueError("Only .pptx templates are supported.")

    row = (
        db.query(DocumentTemplate)
        .filter(
            DocumentTemplate.user_id == user_id,
            DocumentTemplate.public_id == template_id,
            DocumentTemplate.is_active == True,  # noqa: E712
        )
        .first()
    )
    if not row:
        return None

    # Validate before overwriting the existing file.
    prs = Presentation(BytesIO(data))
    logo_asset = None
    try:
        logo_asset = _extract_logo_asset(prs)
    except Exception:
        logger.exception("Failed to extract logo asset from replacement template %s", template_id)

    path = template_path_for_row(row)
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("wb") as f:
        f.write(data)

    now = datetime.now(timezone.utc)
    row.original_filename = filename[:255]
    row.content_type = (content_type or "")[:120] or None
    row.file_size = len(data)
    row.updated_at = now
    db.commit()
    db.refresh(row)

    # template_grammar_for_path is cached by (path, mtime, size), so the new
    # mtime/size from the overwrite above naturally bypasses the stale cache
    # entry for the old content.
    try:
        grammar = template_grammar_for_path(path)
        brand_profile = brand_profile_from_template_grammar(
            grammar, user_id=user_id, profile_id=f"template:{row.public_id}"
        )
        brand_profile.source_template_id = row.public_id
        if logo_asset is not None:
            brand_profile.logo_assets = [logo_asset]
        design_system_id = design_system_id_for_template(user_id, row.public_id)
        write_brand_design_system(brand_profile, design_system_id=design_system_id)
        row.design_system_id = design_system_id
        row.updated_at = datetime.now(timezone.utc)
        db.commit()
        db.refresh(row)
    except Exception:
        logger.exception("Failed to regenerate brand design_system for replaced template %s", row.public_id)

    return row


def archive_user_template(db, user_id: str, template_id: str) -> bool:
    row = (
        db.query(DocumentTemplate)
        .filter(
            DocumentTemplate.user_id == user_id,
            DocumentTemplate.public_id == template_id,
            DocumentTemplate.is_active == True,  # noqa: E712
        )
        .first()
    )
    if not row:
        return False
    row.is_active = False
    row.updated_at = datetime.now(timezone.utc)
    db.commit()
    path = template_path_for_row(row)
    try:
        if path.exists():
            archive_dir = path.parent / ".archived"
            archive_dir.mkdir(exist_ok=True)
            shutil.move(str(path), str(archive_dir / path.name))
    except Exception:
        # DB archive is authoritative; filesystem cleanup can be retried later.
        pass
    return True
