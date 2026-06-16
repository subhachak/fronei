from __future__ import annotations

import base64
from io import BytesIO
from typing import Any


def inspect_pptx_base64(pptx_base64: str) -> list[dict[str, Any]]:
    """Inspect PPTX structure for obvious render risks.

    This is intentionally deterministic and cheap. It does not replace the
    existing LibreOffice/vision QA, but it catches empty slides and excessive
    text density before heavier rendering kicks in.
    """

    if not pptx_base64:
        return [{"type": "missing_artifact", "message": "No PPTX bytes were produced."}]
    try:
        from pptx import Presentation

        prs = Presentation(BytesIO(base64.b64decode(pptx_base64)))
    except Exception as exc:
        return [{"type": "pptx_parse_failed", "message": str(exc)}]

    issues: list[dict[str, Any]] = []
    for index, slide in enumerate(prs.slides, 1):
        text = "\n".join(
            getattr(shape, "text", "") or ""
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
        ).strip()
        if not text:
            issues.append({"type": "empty_slide", "slide": index, "message": "Slide has no text content."})
        if len(text) > 1200:
            issues.append({"type": "dense_slide", "slide": index, "message": "Slide text may be too dense."})
    return issues
