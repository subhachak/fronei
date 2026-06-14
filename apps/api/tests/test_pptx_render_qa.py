import json

import pytest
from pptx import Presentation
from pptx.util import Inches

from app.services.document_generator import generate_pptx_bytes
from app.services.pptx_render_qa import render_qa_available, run_pptx_render_qa

pytestmark = pytest.mark.skipif(
    not render_qa_available(), reason="LibreOffice/poppler not installed in this environment"
)


def _pptx_bytes_from_presentation(prs: Presentation) -> bytes:
    from io import BytesIO

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_render_qa_reports_clean_deck_with_no_issues():
    content = json.dumps({
        "title": "Quarterly Update",
        "slides": [
            {"layout": "bullets", "title": "Highlights", "bullets": ["Revenue up 12%", "Churn down 2pts"]},
        ],
    })

    pptx_bytes = generate_pptx_bytes("Quarterly Update", content)
    result = run_pptx_render_qa(pptx_bytes)

    assert result["available"] is True
    assert result["slide_count"] >= 2  # title slide + content slide
    assert isinstance(result["issues"], list)
    # A short, normal deck shouldn't trip the density/blank heuristics.
    assert result["issues"] == []


def test_render_qa_flags_blank_slide():
    prs = Presentation()
    # "Blank" layout (index 6 in the standard master) with no shapes added.
    prs.slides.add_slide(prs.slide_layouts[6])

    pptx_bytes = _pptx_bytes_from_presentation(prs)
    result = run_pptx_render_qa(pptx_bytes)

    assert result["available"] is True
    assert result["slide_count"] == 1
    assert any(issue["type"] == "blank" for issue in result["issues"])


def test_render_qa_flags_dense_text_slide():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    slide.shapes.title.text = "Wall of text"

    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf = box.text_frame
    tf.word_wrap = True
    paragraph_text = "This sentence is repeated many times to simulate an overflowing slide. "
    tf.text = paragraph_text * 60  # ~4300 characters, well above the density threshold

    pptx_bytes = _pptx_bytes_from_presentation(prs)
    result = run_pptx_render_qa(pptx_bytes)

    assert result["available"] is True
    assert any(issue["type"] == "dense_text" for issue in result["issues"])


def test_render_qa_unavailable_does_not_raise(monkeypatch):
    monkeypatch.setattr("app.services.pptx_render_qa.SOFFICE_BIN", None)

    content = json.dumps({
        "title": "Quarterly Update",
        "slides": [{"layout": "bullets", "title": "Highlights", "bullets": ["Revenue up 12%"]}],
    })
    pptx_bytes = generate_pptx_bytes("Quarterly Update", content)

    result = run_pptx_render_qa(pptx_bytes)

    assert result["available"] is False
    assert result["issues"] == []
    assert "reason" in result
